"""Offline Beamer-to-HTML rendering and chapter finalization."""

from __future__ import annotations

import base64
import hashlib
import html
import json
import math
import re
import shutil
import tempfile
import textwrap
import unicodedata
from collections.abc import Iterable, Sequence
from dataclasses import asdict, dataclass, field, replace
from pathlib import Path
from typing import Any

from pylatexenc.latex2text import LatexNodes2Text

from src.beamer_preflight import normalize_beamer_file
from src.compile import LaTeXCompiler
from src.slide_images import (
    IMAGE_PIPELINE_VERSION,
    ChapterImageResult,
    ImageGenerationConfig,
    append_image_statistics,
    augment_deck_with_generated_images,
    commit_image_result,
    desired_image_state,
    discard_image_result,
    image_cache_is_current,
    image_request_fingerprint,
    load_image_generation_config,
)
from src.slide_style import (
    CourseSlideStyle,
    FONT_FAMILIES,
    FrontendSlidesAssets,
    FrontendSlidesError,
    STYLE_FILENAME,
    ensure_course_slide_style,
    load_assets,
    load_course_slide_style,
    presentation_design_markdown,
    renderer_theme,
    sha256_file,
    slide_gen_asset_root,
    write_presentation_design_result,
)

__all__ = [
    "ChapterFrontendResult",
    "CourseSlideStyle",
    "FrontendSlidesError",
    "ensure_course_slide_style",
    "finalize_chapter",
    "load_course_slide_style",
    "presentation_design_markdown",
    "write_presentation_design_result",
]


# ---------------------------------------------------------------------------
# Data models
# ---------------------------------------------------------------------------

@dataclass
class ContentElement:
    kind: str
    title: str | None = None
    text: str = ""
    items: list["ListItem"] = field(default_factory=list)
    rows: list[list[str]] = field(default_factory=list)
    ordered: bool = False
    start: int = 1
    raw: str = ""
    language: str | None = None
    children: list["ContentElement"] = field(default_factory=list)
    image_data_uri: str | None = None
    image_labels: list[str] = field(default_factory=list)
    env: str | None = None


@dataclass
class ListItem:
    text: str
    children: list[ContentElement] = field(default_factory=list)


@dataclass
class BeamerSlide:
    index: int
    title: str
    elements: list[ContentElement]
    raw_tex: str
    is_titlepage: bool = False
    frame_title_argument: str | None = None


@dataclass
class BeamerDeck:
    source_path: Path
    title: str
    slides: list[BeamerSlide]
    metadata: dict[str, str] = field(default_factory=dict)
    unsupported_environments: list[str] = field(default_factory=list)

    @property
    def slide_count(self) -> int:
        return len(self.slides)


@dataclass
class StyleCandidate:
    id: str
    name: str
    source: str
    preset_name: str | None = None
    slug: str | None = None






@dataclass(frozen=True)
class ChapterFrontendResult:
    html_path: Path
    latex_pdf_path: Path
    html_pdf_path: Path
    html_pptx_path: Path
    manifest_path: Path
    slide_count: int
    skipped: bool = False




# ---------------------------------------------------------------------------
# Beamer parsing
# ---------------------------------------------------------------------------

FRAME_RE = re.compile(r"\\begin\{frame\}|\\frame(?![A-Za-z@])")


ENV_RE = re.compile(r"\\begin\{([A-Za-z*]+)\}")


ITEM_TOKEN_RE = re.compile(r"\\(begin|end)\{([A-Za-z*]+)\}|\\item(?:\[[^\]]*\])?")


TRANSPARENT_ENVS = {
    "center",
    "centering",
    "columns",
    "column",
    "minipage",
    "adjustbox",
    "small",
    "footnotesize",
    "scriptsize",
    "normalsize",
    "large",
    "Large",
}


MATH_ENVS = {
    "equation",
    "equation*",
    "align",
    "align*",
    "alignat",
    "alignat*",
    "gather",
    "gather*",
    "eqnarray",
    "eqnarray*",
    "multline",
    "multline*",
}


ENV_REQUIRED_ARG_COUNTS = {
    "minipage": 1,
    "column": 1,
    "adjustbox": 1,
    "alignat": 1,
    "alignat*": 1,
    "tabular": 1,
}


RESIZEBOX_RE = re.compile(r"\\resizebox\*?\s*(?=\{)")


MATH_SPAN_RE = re.compile(
    r"\\\[(?P<display>.+?)\\\]"
    r"|\\\((?P<inline>.+?)\\\)"
    r"|(?<!\\)\$\$(?P<dollars2>[^$]+?)\$\$"
    r"|(?<!\\)\$(?P<dollars>[^$]+?)(?<!\\)\$",
    re.DOTALL,
)


_MATH_TOKEN = "MATHSPAN{index}ENDSPAN"


_MATH_TOKEN_RE = re.compile(r"MATHSPAN(\d+)ENDSPAN")


def parse_beamer(path: Path | str) -> BeamerDeck:
    source_path = Path(path)
    source = source_path.read_text(encoding="utf-8")
    return parse_beamer_source(source, source_path=source_path)


def parse_beamer_source(
    source: str,
    *,
    source_path: Path | str = Path("slides.tex"),
) -> BeamerDeck:
    """Parse Beamer source already held in memory.

    Chapter generation uses this entry point to compile speaker-note markers
    against the exact finalized source before ``slides.tex`` is written.
    """
    source_path = Path(source_path)
    metadata = _parse_metadata(source)
    document = _document_body(source)
    slides: list[BeamerSlide] = []
    unsupported: list[str] = []

    pos = 0
    while True:
        match = FRAME_RE.search(document, pos)
        if not match:
            break
        if match.group(0).startswith(r"\begin"):
            slide, pos = _parse_frame_environment(document, match.start(), len(slides) + 1)
        else:
            slide, pos = _parse_frame_command(document, match.start(), len(slides) + 1)
        slides.append(slide)
        unsupported.extend(_find_unsupported(slide.elements))

    title = _clean_text(metadata.get("title", "")) or _infer_title_from_slides(slides) or source_path.stem
    return BeamerDeck(
        source_path=source_path,
        title=title,
        slides=slides,
        metadata=metadata,
        unsupported_environments=sorted(set(unsupported)),
    )






def _parse_metadata(source: str) -> dict[str, str]:
    metadata: dict[str, str] = {}
    for key in ("title", "subtitle", "author", "institute", "date"):
        match = re.search(rf"\\{key}(?:\[[^\]]*\])?\s*\{{", source)
        if not match:
            continue
        group = _parse_group(source, match.end() - 1, "{", "}")
        if group:
            metadata[key] = _clean_text(group[0])
    return metadata


def _document_body(source: str) -> str:
    begin = source.find(r"\begin{document}")
    end = source.rfind(r"\end{document}")
    if begin == -1:
        return source
    begin += len(r"\begin{document}")
    if end == -1 or end <= begin:
        return source[begin:]
    return source[begin:end]


def _parse_frame_environment(source: str, start: int, index: int) -> tuple[BeamerSlide, int]:
    command_end = start + len(r"\begin{frame}")
    cursor = _skip_ws(source, command_end)
    while cursor < len(source) and source[cursor] == "[":
        group = _parse_group(source, cursor, "[", "]")
        if not group:
            break
        cursor = _skip_ws(source, group[1])

    frame_title_argument = None
    if cursor < len(source) and source[cursor] == "{":
        group = _parse_group(source, cursor, "{", "}")
        if group:
            frame_title_argument = _clean_text(group[0])
            cursor = _skip_ws(source, group[1])

    end_start = source.find(r"\end{frame}", cursor)
    if end_start == -1:
        end_start = len(source)
        end_cursor = len(source)
    else:
        end_cursor = end_start + len(r"\end{frame}")

    body = source[cursor:end_start]
    explicit_title, body_without_title = _extract_frame_title(body)
    title = explicit_title or frame_title_argument or ""
    is_titlepage = r"\titlepage" in body_without_title
    elements = [] if is_titlepage else _parse_elements(body_without_title)
    slide = BeamerSlide(
        index=index,
        title=title,
        elements=elements,
        raw_tex=source[start:end_cursor],
        is_titlepage=is_titlepage,
        frame_title_argument=frame_title_argument,
    )
    return slide, end_cursor


def _parse_frame_command(source: str, start: int, index: int) -> tuple[BeamerSlide, int]:
    cursor = _skip_ws(source, start + len(r"\frame"))
    while cursor < len(source) and source[cursor] == "[":
        group = _parse_group(source, cursor, "[", "]")
        if not group:
            break
        cursor = _skip_ws(source, group[1])
    group = _parse_group(source, cursor, "{", "}")
    if not group:
        body = ""
        end_cursor = cursor
    else:
        body, end_cursor = group
    explicit_title, body_without_title = _extract_frame_title(body)
    is_titlepage = r"\titlepage" in body_without_title
    slide = BeamerSlide(
        index=index,
        title=explicit_title,
        elements=[] if is_titlepage else _parse_elements(body_without_title),
        raw_tex=source[start:end_cursor],
        is_titlepage=is_titlepage,
    )
    return slide, end_cursor


def _extract_frame_title(body: str) -> tuple[str, str]:
    match = re.search(r"\\frametitle\s*\{", body)
    if not match:
        return "", body
    group = _parse_group(body, match.end() - 1, "{", "}")
    if not group:
        return "", body
    title, end = group
    return _clean_text(title), body[: match.start()] + body[end:]


def _parse_elements(source: str) -> list[ContentElement]:
    source = _strip_comments(source)
    source = _strip_resizebox(source)
    elements: list[ContentElement] = []
    pos = 0
    while pos < len(source):
        match = ENV_RE.search(source, pos)
        if not match:
            _append_text_element(elements, source[pos:])
            break

        _append_text_element(elements, source[pos : match.start()])
        env = match.group(1)
        cursor = match.end()
        option_text = None
        if cursor < len(source) and source[cursor] == "[":
            option = _parse_group(source, cursor, "[", "]")
            if option:
                option_text, cursor = option

        block_title = None
        if env == "block":
            cursor = _skip_ws(source, cursor)
            if cursor < len(source) and source[cursor] == "{":
                group = _parse_group(source, cursor, "{", "}")
                if group:
                    block_title, cursor = _clean_text(group[0]), group[1]

        for _ in range(ENV_REQUIRED_ARG_COUNTS.get(env, 0)):
            cursor = _skip_ws(source, cursor)
            group = _parse_group(source, cursor, "{", "}")
            if not group:
                break
            cursor = group[1]

        end_range = _find_matching_environment(source, env, cursor)
        if not end_range:
            _append_text_element(elements, source[match.start() :])
            break
        inner = source[cursor : end_range[0]]
        pos = end_range[1]

        if env == "block":
            elements.append(
                ContentElement(kind="block", title=block_title, children=_parse_elements(inner), raw=inner)
            )
        elif env in {"itemize", "enumerate"}:
            elements.append(
                ContentElement(
                    kind="list",
                    ordered=env == "enumerate",
                    items=_parse_list_items(inner),
                    raw=inner,
                )
            )
        elif env in MATH_ENVS:
            elements.append(
                ContentElement(kind="equation", text=_strip_math_labels(inner).strip(), raw=inner, env=env)
            )
        elif env in {"table", "tabular"}:
            elements.append(ContentElement(kind="table", rows=_parse_table_rows(inner), raw=inner, env=env))
        elif env == "lstlisting":
            elements.append(
                ContentElement(
                    kind="code",
                    text=textwrap.dedent(inner).strip("\n"),
                    raw=inner,
                    language=_language_from_option(option_text),
                )
            )
        elif env == "verbatim":
            elements.append(
                ContentElement(kind="code", text=textwrap.dedent(inner).strip("\n"), raw=inner, language=None)
            )
        elif env in TRANSPARENT_ENVS:
            elements.extend(_parse_elements(inner))
        elif env not in {"tikzpicture", "axis"} and r"\begin{tabular}" in inner:
            # Safety net: never let a real table reach the raw-LaTeX fallback.
            elements.extend(_parse_elements(inner))
        else:
            elements.append(ContentElement(kind="raw", title=env, raw=inner.strip()))
    return [element for element in elements if _element_has_content(element)]


def _parse_list_items(source: str) -> list[ListItem]:
    spans: list[tuple[int, int]] = []
    depth = 0
    for token in ITEM_TOKEN_RE.finditer(source):
        if token.group(0).startswith(r"\begin"):
            depth += 1
        elif token.group(0).startswith(r"\end"):
            depth = max(0, depth - 1)
        elif depth == 0:
            spans.append((token.start(), token.end()))
    items: list[ListItem] = []
    for i, (_, start) in enumerate(spans):
        end = spans[i + 1][0] if i + 1 < len(spans) else len(source)
        segment = source[start:end]
        nested = ENV_RE.search(segment)
        if nested:
            text_part = segment[: nested.start()]
            child_part = segment[nested.start() :]
            children = _parse_elements(child_part)
        else:
            text_part = segment
            children = []
        text = _clean_text(text_part)
        if text or children:
            items.append(ListItem(text=text, children=children))
    return items


def _parse_table_rows(source: str) -> list[list[str]]:
    tabular = re.search(r"\\begin\{tabular\}(?:\{[^}]*\})?", source)
    if tabular:
        end = source.find(r"\end{tabular}", tabular.end())
        if end != -1:
            source = source[tabular.end() : end]
    source = re.sub(r"\\(hline|toprule|midrule|bottomrule)\b", "", source)
    source = re.sub(r"\\cline\s*\{[^}]*\}", "", source)
    # Flatten \multicolumn{n}{fmt}{body} to its body; the braces around the
    # body are removed later by _clean_text.
    source = re.sub(r"\\multicolumn\s*\{[^}]*\}\s*\{[^}]*\}", "", source)
    rows: list[list[str]] = []
    for raw_row in re.split(r"\\\\", source):
        # Keep empty cells so columns stay aligned; drop rows that are entirely empty.
        cells = [_clean_text(cell) for cell in re.split(r"(?<!\\)&", raw_row)]
        if any(cells):
            rows.append(cells)
    return rows


def _append_text_element(elements: list[ContentElement], source: str) -> None:
    text = _clean_text(source)
    if text:
        elements.append(ContentElement(kind="text", text=text, raw=source))


def _element_has_content(element: ContentElement) -> bool:
    return bool(element.text or element.items or element.rows or element.raw or element.children)


def _find_matching_environment(source: str, env: str, start: int) -> tuple[int, int] | None:
    pattern = re.compile(rf"\\(begin|end)\{{{re.escape(env)}\}}")
    depth = 1
    for match in pattern.finditer(source, start):
        if match.group(1) == "begin":
            depth += 1
        else:
            depth -= 1
            if depth == 0:
                return match.start(), match.end()
    return None


def _parse_group(source: str, start: int, opener: str, closer: str) -> tuple[str, int] | None:
    if start >= len(source) or source[start] != opener:
        return None
    depth = 1
    cursor = start + 1
    while cursor < len(source):
        char = source[cursor]
        if char == "\\":
            cursor += 2
            continue
        if char == opener:
            depth += 1
        elif char == closer:
            depth -= 1
            if depth == 0:
                return source[start + 1 : cursor], cursor + 1
        cursor += 1
    return None


def _skip_ws(source: str, cursor: int) -> int:
    while cursor < len(source) and source[cursor].isspace():
        cursor += 1
    return cursor


def _strip_resizebox(source: str) -> str:
    """Rewrite \\resizebox{w}{h}{content} to just content (nested braces safe)."""
    while True:
        match = RESIZEBOX_RE.search(source)
        if not match:
            return source
        cursor = match.end()
        groups: list[str] = []
        for _ in range(3):
            cursor = _skip_ws(source, cursor)
            group = _parse_group(source, cursor, "{", "}")
            if not group:
                break
            groups.append(group[0])
            cursor = group[1]
        if len(groups) != 3:
            # Malformed command: drop the token itself so the loop terminates.
            source = source[: match.start()] + source[match.end() :]
            continue
        source = source[: match.start()] + groups[2] + source[cursor:]


def _strip_math_labels(source: str) -> str:
    return re.sub(r"\\label\s*\{[^}]*\}", "", source)


_DISPLAY_SAFE_ENV_MAP = {
    "align": "aligned",
    "align*": "aligned",
    "alignat": "alignedat",
    "alignat*": "alignedat",
    "flalign": "aligned",
    "flalign*": "aligned",
    "eqnarray": "aligned",
    "eqnarray*": "aligned",
    "gather": "gathered",
    "gather*": "gathered",
    "multline": "gathered",
    "multline*": "gathered",
    "equation": "",
    "equation*": "",
}


_ENV_TOKEN_RE = re.compile(r"\\(begin|end)\{([A-Za-z*]+)\}")


def normalize_display_math(body: str) -> str:
    """Rewrite environments nested in display math to forms MathJax accepts inside \\[...\\]."""

    def swap(match: re.Match[str]) -> str:
        replacement = _DISPLAY_SAFE_ENV_MAP.get(match.group(2))
        if replacement is None:
            return match.group(0)
        if not replacement:
            return ""
        return f"\\{match.group(1)}{{{replacement}}}"

    return _ENV_TOKEN_RE.sub(swap, body)


def _strip_comments(source: str) -> str:
    lines = []
    for line in source.splitlines():
        lines.append(re.sub(r"(?<!\\)%.*$", "", line))
    return "\n".join(lines)


def _clean_text(source: str) -> str:
    if not source:
        return ""
    source = _strip_comments(source)
    source, math_spans = _protect_math(source)
    source = re.sub(r"\\setcounter\{[^}]+\}\{[^}]+\}", "", source)
    source = source.replace("~", " ")
    source = re.sub(r"\\separator\b", " ", source)
    try:
        text = LatexNodes2Text().latex_to_text(source)
    except Exception:
        text = source
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = _restore_math(text, math_spans)
    return text.strip()


def _protect_math(source: str) -> tuple[str, list[str]]:
    """Replace math spans with plain-text tokens that survive latex-to-text conversion.

    Spans are normalized to MathJax delimiters: \\( \\) for inline, \\[ \\] for display.
    """
    spans: list[str] = []

    def stash(match: re.Match[str]) -> str:
        display = match.group("display") or match.group("dollars2")
        inline = match.group("inline") or match.group("dollars")
        if display is not None:
            spans.append(f"\\[{normalize_display_math(display.strip())}\\]")
        else:
            spans.append(f"\\({inline.strip()}\\)")
        return _MATH_TOKEN.format(index=len(spans) - 1)

    return MATH_SPAN_RE.sub(stash, source), spans


def _restore_math(text: str, spans: list[str]) -> str:
    if not spans:
        return text

    def unstash(match: re.Match[str]) -> str:
        index = int(match.group(1))
        return spans[index] if index < len(spans) else match.group(0)

    return _MATH_TOKEN_RE.sub(unstash, text)


def _language_from_option(option_text: str | None) -> str | None:
    if not option_text:
        return None
    match = re.search(r"language\s*=\s*([A-Za-z0-9_+-]+)", option_text)
    return match.group(1) if match else None


def _find_unsupported(elements: list[ContentElement]) -> list[str]:
    unsupported: list[str] = []
    for element in elements:
        if element.kind == "raw" and element.title:
            unsupported.append(element.title)
        unsupported.extend(_find_unsupported(element.children))
        for item in element.items:
            unsupported.extend(_find_unsupported(item.children))
    return unsupported


def _infer_title_from_slides(slides: list[BeamerSlide]) -> str:
    for slide in slides:
        if slide.is_titlepage:
            continue
        title = slide.title.strip()
        if not title:
            continue
        title = re.sub(r"^Introduction to\s+", "", title, flags=re.IGNORECASE)
        title = re.split(r"\s+-\s+", title)[0].strip()
        return title
    return ""




# ---------------------------------------------------------------------------
# Speaker notes
# ---------------------------------------------------------------------------

NOTE_MARKER_RE = re.compile(
    r"^<!--\s*speaker-note:\s*(slide-(?P<number>\d{3,}))\s*-->\s*$",
    re.MULTILINE,
)


NOTE_HEADING_RE = re.compile(
    r"^##\s+Slide\s+(?P<number>\d{3,}):\s*(?P<title>.+?)\s*$",
    re.MULTILINE,
)


LEGACY_HEADING_RE = re.compile(
    r"^##\s+Section\s+(?P<number>\d+):\s*(?P<title>.+?)\s*$",
    re.MULTILINE,
)


TRAILING_SEPARATOR_RE = re.compile(
    r"\n[ \t]*---[ \t]*(?:\n[ \t]*)*\Z"
)


@dataclass(frozen=True)
class SpeakerNote:
    """One validated note correlated to one rendered slide."""

    id: str
    index: int
    title: str
    text: str

    @property
    def sha256(self) -> str:
        return hashlib.sha256(self.text.encode("utf-8")).hexdigest()

    def to_payload(self) -> dict[str, object]:
        return {
            "id": self.id,
            "index": self.index,
            "title": self.title,
            "text": self.text,
        }

    def to_manifest_entry(self) -> dict[str, object]:
        return {
            "id": self.id,
            "index": self.index,
            "title": self.title,
            "note_sha256": self.sha256,
        }


def stable_slide_id(index: int) -> str:
    if index < 1:
        raise ValueError("Slide indexes are one-based.")
    return f"slide-{index:03d}"


def slide_title(deck: BeamerDeck, slide: BeamerSlide) -> str:
    if slide.is_titlepage:
        return deck.title
    return slide.title.strip() or "Untitled"


def normalize_title(value: str) -> str:
    """Normalize titles for correlation without discarding meaningful words."""
    normalized = unicodedata.normalize("NFKC", value)
    normalized = re.sub(r"[*_`]+", "", normalized)
    normalized = re.sub(r"\s+", " ", normalized).strip()
    return normalized.casefold()


def deterministic_title_note(title: str) -> str:
    return (
        f"Welcome learners to {title}. Briefly orient them to the chapter's "
        "purpose, learning goals, and the progression of the slides that follow."
    )


def parse_speaker_notes(markdown: str) -> list[SpeakerNote]:
    """Parse canonical note blocks from ``script.md``.

    Each block begins with a stable HTML-comment marker followed by a visible
    heading. Requiring both keeps the document pleasant to read while giving
    the renderer an unambiguous machine contract.
    """
    markers = list(NOTE_MARKER_RE.finditer(markdown))
    if not markers:
        raise FrontendSlidesError(
            "script.md has no speaker-note markers; expected "
            "'<!-- speaker-note: slide-001 -->'."
        )
    marker_ids = [marker.group(1) for marker in markers]
    duplicate_ids = sorted(
        note_id for note_id in set(marker_ids) if marker_ids.count(note_id) > 1
    )
    if duplicate_ids:
        raise FrontendSlidesError(
            "Duplicate speaker-note ID: " + ", ".join(duplicate_ids) + "."
        )

    notes: list[SpeakerNote] = []
    for position, marker in enumerate(markers):
        block_start = marker.end()
        block_end = markers[position + 1].start() if position + 1 < len(markers) else len(markdown)
        block = markdown[block_start:block_end]
        heading = NOTE_HEADING_RE.search(block)
        if not heading:
            raise FrontendSlidesError(
                f"{marker.group(1)} is missing its '## Slide NNN: Title' heading."
            )
        prefix = block[: heading.start()]
        if prefix.strip():
            raise FrontendSlidesError(
                f"{marker.group(1)} has unexpected content before its slide heading."
            )
        marker_number = int(marker.group("number"))
        heading_number = int(heading.group("number"))
        if marker_number != heading_number:
            raise FrontendSlidesError(
                f"{marker.group(1)} disagrees with heading number {heading_number:03d}."
            )
        body = block[heading.end():]
        separator = TRAILING_SEPARATOR_RE.search(body)
        if separator:
            body = body[: separator.start()]
        text = body.strip()
        notes.append(
            SpeakerNote(
                id=marker.group(1),
                index=marker_number,
                title=heading.group("title").strip(),
                text=text,
            )
        )
    return notes


def load_speaker_notes(path: Path | str) -> list[SpeakerNote]:
    script_path = Path(path)
    if not script_path.is_file() or script_path.stat().st_size == 0:
        raise FrontendSlidesError(f"Speaker notes are missing or empty: {script_path}")
    return parse_speaker_notes(script_path.read_text(encoding="utf-8"))


def validate_speaker_notes(
    deck: BeamerDeck,
    notes: Sequence[SpeakerNote],
) -> list[str]:
    errors: list[str] = []
    if len(notes) != deck.slide_count:
        errors.append(
            f"Speaker-note count {len(notes)} does not match slide count {deck.slide_count}."
        )

    seen: set[str] = set()
    for expected_index, note in enumerate(notes, 1):
        expected_id = stable_slide_id(expected_index)
        if note.id in seen:
            errors.append(f"Duplicate speaker-note ID: {note.id}.")
        seen.add(note.id)
        if note.index != expected_index or note.id != expected_id:
            errors.append(
                f"Speaker note {expected_index} must use ordered ID {expected_id}, "
                f"found {note.id}."
            )
        if not note.text.strip():
            errors.append(f"{note.id} has empty speaker notes.")
        if expected_index <= deck.slide_count:
            expected_title = slide_title(deck, deck.slides[expected_index - 1])
            if normalize_title(note.title) != normalize_title(expected_title):
                errors.append(
                    f"{note.id} title {note.title!r} does not match rendered "
                    f"slide title {expected_title!r}."
                )
    return errors


def correlate_speaker_notes(
    deck: BeamerDeck,
    markdown_or_notes: str | Sequence[SpeakerNote],
) -> list[SpeakerNote]:
    notes = (
        parse_speaker_notes(markdown_or_notes)
        if isinstance(markdown_or_notes, str)
        else list(markdown_or_notes)
    )
    errors = validate_speaker_notes(deck, notes)
    if errors:
        raise FrontendSlidesError("Speaker-note correlation failed: " + "; ".join(errors))
    return notes


def render_speaker_notes_markdown(
    deck: BeamerDeck,
    content_notes: Sequence[str],
    *,
    document_title: str = "Slides Script",
) -> str:
    """Create canonical ``script.md`` from finalized slide order.

    ``content_notes`` excludes title pages. Title-page notes are deterministic,
    so model-generated notes remain mapped only to the content frames they were
    authored from.
    """
    expected_content = sum(not slide.is_titlepage for slide in deck.slides)
    if len(content_notes) != expected_content:
        raise FrontendSlidesError(
            f"Cannot compile script.md: received {len(content_notes)} content notes "
            f"for {expected_content} content slides."
        )
    content_iterator = iter(content_notes)
    notes: list[SpeakerNote] = []
    for index, slide in enumerate(deck.slides, 1):
        title = slide_title(deck, slide)
        text = (
            deterministic_title_note(title)
            if slide.is_titlepage
            else str(next(content_iterator)).strip()
        )
        notes.append(
            SpeakerNote(
                id=stable_slide_id(index),
                index=index,
                title=title,
                text=text,
            )
        )
    errors = validate_speaker_notes(deck, notes)
    if errors:
        raise FrontendSlidesError("Cannot compile script.md: " + "; ".join(errors))
    return _render_markdown(document_title, notes)


def upgrade_legacy_script(deck: BeamerDeck, markdown: str) -> str:
    """Upgrade a legacy ``## Section`` script when titles align exactly.

    A mismatch is intentionally fatal: silently attaching a plausible note to
    the wrong slide is worse than requiring note-only repair.
    """
    if NOTE_MARKER_RE.search(markdown):
        correlate_speaker_notes(deck, markdown)
        return markdown

    headings = list(LEGACY_HEADING_RE.finditer(markdown))
    expected_slides = [slide for slide in deck.slides if not slide.is_titlepage]
    if len(headings) != len(expected_slides):
        raise FrontendSlidesError(
            "Legacy script cannot be upgraded safely: "
            f"{len(headings)} sections for {len(expected_slides)} content slides."
        )

    content_notes: list[str] = []
    for position, (heading, slide) in enumerate(zip(headings, expected_slides)):
        expected_title = slide_title(deck, slide)
        actual_title = heading.group("title").strip()
        if normalize_title(actual_title) != normalize_title(expected_title):
            raise FrontendSlidesError(
                "Legacy script cannot be upgraded safely: "
                f"section {position + 1} title {actual_title!r} does not match "
                f"rendered slide title {expected_title!r}."
            )
        block_end = headings[position + 1].start() if position + 1 < len(headings) else len(markdown)
        body = markdown[heading.end():block_end]
        separator = TRAILING_SEPARATOR_RE.search(body)
        if separator:
            body = body[: separator.start()]
        note = body.strip()
        if not note:
            raise FrontendSlidesError(
                f"Legacy script cannot be upgraded safely: section {position + 1} is empty."
            )
        content_notes.append(note)

    heading = markdown.splitlines()[0].lstrip("# ").strip() if markdown.strip() else "Slides Script"
    return render_speaker_notes_markdown(
        deck,
        content_notes,
        document_title=heading or "Slides Script",
    )


def repair_speaker_notes_markdown(
    deck: BeamerDeck,
    *,
    document_title: str = "Slides Script",
) -> str:
    """Create correlated replacement notes from the exact parsed slide content.

    This is intentionally a note-only recovery path for legacy decks whose
    historical scripts cannot be aligned safely. It never changes slides.tex
    or assessment.md.
    """
    content_notes = [
        _deterministic_content_note(slide)
        for slide in deck.slides
        if not slide.is_titlepage
    ]
    return render_speaker_notes_markdown(
        deck,
        content_notes,
        document_title=document_title,
    )


def notes_manifest(notes: Iterable[SpeakerNote]) -> list[dict[str, object]]:
    return [note.to_manifest_entry() for note in notes]


def _render_markdown(document_title: str, notes: Sequence[SpeakerNote]) -> str:
    blocks = [f"# {document_title.strip().lstrip('#').strip()}\n"]
    for note in notes:
        blocks.append(
            f"<!-- speaker-note: {note.id} -->\n"
            f"## Slide {note.index:03d}: {note.title}\n\n"
            f"{note.text.strip()}\n\n"
            "---\n"
        )
    return "\n".join(blocks).rstrip() + "\n"


def _deterministic_content_note(slide: BeamerSlide) -> str:
    visible = " ".join(
        fragment
        for element in slide.elements
        for fragment in _element_fragments(element)
        if fragment
    )
    visible = re.sub(r"\s+", " ", visible).strip()
    if len(visible) > 1800:
        visible = visible[:1797].rstrip() + "..."
    detail = (
        f" Walk learners through these visible elements in order: {visible}"
        if visible
        else ""
    )
    return (
        f"Present {slide.title.strip() or 'this slide'}.{detail} "
        "Explain the relationships and any formulas directly from the slide, "
        "then check understanding before advancing."
    )


def _element_fragments(element: ContentElement) -> list[str]:
    fragments = [element.title or "", element.text or ""]
    for item in element.items:
        fragments.extend(_item_fragments(item))
    for row in element.rows:
        fragments.extend(str(cell) for cell in row)
    for child in element.children:
        fragments.extend(_element_fragments(child))
    if not any(fragment.strip() for fragment in fragments) and element.raw:
        fragments.append(element.raw)
    return [str(fragment).strip() for fragment in fragments if str(fragment).strip()]


def _item_fragments(item: ListItem) -> list[str]:
    fragments = [item.text]
    for child in item.children:
        fragments.extend(_element_fragments(child))
    return [str(fragment).strip() for fragment in fragments if str(fragment).strip()]


# ---------------------------------------------------------------------------
# Layout weights
# ---------------------------------------------------------------------------

_CHARS_PER_LINE = 90


_DISPLAY_MATH_RE = re.compile(
    r"\\\[.*?\\\]"
    r"|\\begin\{(?:equation|align|gather|multline)\*?\}.*?"
    r"\\end\{(?:equation|align|gather|multline)\*?\}"
    r"|\$\$.*?\$\$",
    re.DOTALL,
)


def element_weight(element: ContentElement) -> int:
    if element.kind == "list":
        return sum(item_weight(item) for item in element.items)
    if element.kind == "block":
        return 1 + sum(element_weight(child) for child in element.children)
    if element.kind == "table":
        return max(3, len(element.rows))
    if element.kind in {"generated_image", "user_image"}:
        return 4
    if element.kind == "code":
        # Code renders as a carbon-style image card (~480px, the default) or a
        # <pre>. Weight for the image case: splitting a touch too eagerly is
        # safe, letting a tall image overflow is not.
        if element.image_data_uri:
            return 6
        return max(3, min(element.text.count("\n") + 2, 6))
    if element.kind in {"equation", "raw"}:
        return 2
    return max(1, math.ceil(len(element.text) / _CHARS_PER_LINE)) + _display_math_count(
        element.text
    )


def item_weight(item: ListItem) -> int:
    lines = 1 + len(item.text) // _CHARS_PER_LINE
    return (
        lines
        + _display_math_count(item.text)
        + sum(element_weight(child) for child in item.children)
    )




def _display_math_count(text: str) -> int:
    """Account for the vertical box MathJax adds for every display equation."""
    return len(_DISPLAY_MATH_RE.findall(text))


# ---------------------------------------------------------------------------
# HTML rendering
# ---------------------------------------------------------------------------

















def render_course_presentation_html(
    deck: BeamerDeck,
    style: CourseSlideStyle,
    assets: FrontendSlidesAssets,
    *,
    font_css: str,
    mathjax_src: str = "assets/mathjax/tex-svg.js",
    speaker_notes: Sequence[SpeakerNote] | None = None,
) -> str:
    """Render a complete offline deck from a validated course-wide style."""
    candidate = StyleCandidate(
        id="course-style",
        name=style.selected_style.name,
        source=style.selected_style.source,
        preset_name=(
            style.selected_style.name
            if style.selected_style.source == "preset"
            else None
        ),
        slug=(
            style.selected_style.key
            if style.selected_style.source == "bold_template"
            else None
        ),
    )
    return render_deck_html(
        deck,
        candidate,
        assets,
        theme_override=renderer_theme(style),
        font_css=font_css,
        mathjax_src=mathjax_src,
        layout_rotation=style.presentation_method.layout_rotation,
        speaker_notes=speaker_notes,
    )


def render_deck_html(
    deck: BeamerDeck,
    candidate: StyleCandidate,
    assets: FrontendSlidesAssets,
    *,
    theme_override: dict[str, str] | None = None,
    font_css: str | None = None,
    mathjax_src: str = "https://cdn.jsdelivr.net/npm/mathjax@3/es5/tex-svg.js",
    layout_rotation: list[str] | None = None,
    speaker_notes: Sequence[SpeakerNote] | None = None,
) -> str:
    theme = theme_override or theme_for_candidate(candidate)
    correlated_notes = (
        correlate_speaker_notes(deck, speaker_notes)
        if speaker_notes is not None
        else []
    )
    slides = "\n".join(
        render_slide(deck, slide, theme, layout_rotation) for slide in deck.slides
    )
    manifest_for_page = {
        "title": deck.title,
        "slideCount": deck.slide_count,
        "style": candidate.name,
        "source": str(deck.source_path),
        "speakerNoteCount": len(correlated_notes),
    }
    notes_payload = _json_for_script(
        [note.to_payload() for note in correlated_notes]
    )
    design_note = f"{candidate.name}; generated/custom style recipe"
    return f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{html.escape(deck.title)}</title>
    {font_css if font_css is not None else _font_links(theme)}
    <script>
        window.MathJax = {{
            tex: {{ inlineMath: [['\\\\(', '\\\\)'], ['$', '$']], displayMath: [['\\\\[', '\\\\]']] }},
            svg: {{ fontCache: 'global' }}
        }};
    </script>
    <script defer src="{html.escape(mathjax_src)}"></script>
    <style>
        /* === CSS CUSTOM PROPERTIES (THEME) === */
        :root {{
            --stage-bg: {theme['stage_bg']};
            --slide-bg: {theme['background']};
            --text-primary: {theme['text']};
            --text-secondary: {theme['muted']};
            --accent: {theme['accent']};
            --accent-2: {theme['accent2']};
            --surface: {theme['surface']};
            --surface-alt: {theme['surface_alt']};
            --border: {theme['border']};
            --font-display: {theme['display_font']};
            --font-body: {theme['body_font']};
            --font-mono: {theme['mono_font']};
            --ease-out-expo: cubic-bezier(0.16, 1, 0.3, 1);
        }}

        /* === RESET === */
        * {{ margin: 0; padding: 0; box-sizing: border-box; }}

{_indent(assets.viewport_css, 8)}

        /* === PRESENTATION THEME === */
        body {{ color: var(--text-primary); font-family: var(--font-body); }}
        .slide {{
            background:
                radial-gradient(circle at 86% 12%, {theme['glow']} 0, transparent 330px),
                linear-gradient(135deg, {theme['background']} 0%, {theme['background_alt']} 100%);
            color: var(--text-primary);
        }}
        .slide::before {{
            content: "";
            position: absolute;
            inset: 0;
            background-image:
                linear-gradient({theme['grid']} 1px, transparent 1px),
                linear-gradient(90deg, {theme['grid']} 1px, transparent 1px);
            background-size: 80px 80px;
            opacity: {theme['grid_opacity']};
            pointer-events: none;
        }}
        .slide-frame {{
            position: absolute;
            inset: 52px 64px;
            display: grid;
            grid-template-rows: 48px minmax(0, 1fr) 40px;
            gap: 34px;
        }}
        .slide-topline,
        .slide-footline {{
            position: relative;
            z-index: 1;
            display: flex;
            justify-content: space-between;
            align-items: center;
            color: var(--text-secondary);
            border-color: var(--border);
            font-family: var(--font-mono);
            font-size: 18px;
            letter-spacing: 0.08em;
            text-transform: uppercase;
        }}
        .slide-topline {{ border-bottom: 1px solid var(--border); padding-bottom: 16px; }}
        .slide-footline {{ border-top: 1px solid var(--border); padding-top: 14px; }}

        /* === SLIDE CONTENT + LAYOUT VARIANTS === */
        .slide-content {{
            --fit-scale: 1;
            --density: {theme.get('course_density', '1')};
            position: relative;
            z-index: 1;
            display: grid;
            min-height: 0;
            overflow: hidden;
        }}
        .slide-content.compact {{ --density: 0.92; }}
        .layout-hero {{ grid-template-columns: 1fr; align-content: end; }}
        .layout-split {{
            grid-template-columns: minmax(0, 0.85fr) minmax(0, 1.15fr);
            gap: 76px;
            align-items: stretch;
        }}
        .layout-split .slide-head {{ align-self: center; }}
        .layout-top,
        .layout-top-cols,
        .layout-math {{
            grid-template-rows: auto minmax(0, 1fr);
            gap: 42px;
            align-content: start;
        }}
        .layout-top .slide-body {{ max-width: 1560px; }}
        .layout-top-cols .slide-body {{ grid-template-columns: repeat(2, minmax(0, 1fr)); gap: 32px 68px; }}
        .layout-top-cols .col {{ align-content: start; }}
        .layout-math .slide-body {{ grid-template-columns: minmax(0, 1.05fr) minmax(0, 0.95fr); gap: 76px; }}

        /* === SLIDE HEADER === */
        .slide-head .kicker {{
            margin-bottom: 18px;
            display: inline-flex;
            align-items: center;
            gap: 14px;
            color: var(--accent);
            font-family: var(--font-mono);
            font-size: 20px;
            letter-spacing: 0.2em;
            text-transform: uppercase;
        }}
        .slide-head .kicker::before {{ content: ""; width: 40px; height: 2px; background: var(--accent); }}
        h1,
        h2 {{
            font-family: var(--font-display);
            letter-spacing: 0;
            color: var(--text-primary);
        }}
        h1 {{ font-size: {theme['title_size']}; line-height: 0.94; }}
        h2 {{ font-size: {theme['heading_size']}; line-height: 1.04; max-width: 1520px; }}
        .layout-split h2 {{ max-width: 640px; }}
        .title-lockup {{
            max-width: 1280px;
            padding-bottom: 40px;
        }}
        .title-subtitle {{
            margin-top: 30px;
            max-width: 940px;
            color: var(--text-secondary);
            font-size: 34px;
            line-height: 1.35;
        }}

        /* === SLIDE BODY (open flow, no nested boxes) === */
        .slide-body {{
            display: grid;
            gap: calc(30px * var(--fit-scale));
            align-content: safe center;
            min-height: 0;
            overflow: hidden;
        }}
        .col {{
            display: grid;
            gap: calc(30px * var(--fit-scale));
            align-content: safe center;
            min-height: 0;
        }}
        p,
        li,
        td,
        th {{
            color: var(--text-primary);
            font-size: calc(26px * var(--density) * var(--fit-scale));
            line-height: 1.45;
        }}
        .text-item {{ color: var(--text-secondary); max-width: 1400px; }}

        /* Open lists with custom accent markers */
        .blist {{ list-style: none; padding: 0; display: grid; gap: calc(18px * var(--fit-scale)); }}
        ol.blist {{ counter-reset: bitem; }}
        ol.blist > li {{ position: relative; padding-left: 2.4em; counter-increment: bitem; }}
        ol.blist > li::before {{
            content: counter(bitem, decimal-leading-zero) ".";
            position: absolute;
            left: 0;
            top: 0;
            color: var(--accent);
            font-family: var(--font-mono);
            font-weight: 700;
        }}
        ul.blist > li {{ position: relative; padding-left: 38px; }}
        ul.blist > li::before {{
            content: "";
            position: absolute;
            left: 2px;
            top: 0.42em;
            width: 12px;
            height: 12px;
            border: 2px solid var(--accent);
            transform: rotate(45deg);
        }}
        li .blist {{ margin-top: calc(12px * var(--fit-scale)); gap: calc(11px * var(--fit-scale)); }}
        li ul.blist {{ padding-left: 4px; }}
        li ul.blist > li {{ padding-left: 30px; }}
        li ul.blist > li::before {{
            left: 4px;
            top: 0.5em;
            width: 8px;
            height: 8px;
            border-color: var(--accent-2);
            border-radius: 2px;
            transform: none;
        }}
        li ol.blist > li {{ padding-left: 1.7em; }}
        li ol.blist > li::before {{ content: counter(bitem) "."; }}
        li li {{ font-size: 0.86em; color: var(--text-secondary); }}
        strong {{ color: var(--accent); font-weight: 800; }}
        mark {{ background: color-mix(in srgb, var(--accent) 35%, transparent); color: inherit; padding: 0 0.15em; }}

        /* Beamer block: accent-edged section, not a box of boxes */
        .content-block {{
            border-left: 4px solid var(--accent);
            padding: 8px 0 8px 30px;
            display: grid;
            gap: calc(18px * var(--fit-scale));
        }}
        .content-block h3 {{
            color: var(--accent);
            font-family: var(--font-display);
            font-size: calc(31px * var(--density) * var(--fit-scale));
            line-height: 1.12;
        }}

        /* Code, tables and raw LaTeX earn a surface; equations stay unboxed. */
        .code-card,
        .table-card,
        .raw-card {{
            background: {theme['panel_fill']};
            border: {theme['panel_border']};
            box-shadow: {theme['shadow']};
            padding: calc(26px * var(--fit-scale)) calc(34px * var(--fit-scale));
            overflow: hidden;
        }}
        .formula {{
            background: transparent;
            border: 0;
            box-shadow: none;
            padding: calc(8px * var(--fit-scale)) calc(4px * var(--fit-scale));
            overflow: visible;
            color: var(--text-primary);
            font-size: calc(27px * var(--density) * var(--fit-scale));
        }}
        .code-card--image {{
            background: transparent;
            border: 0;
            box-shadow: none;
            padding: 0;
        }}
        .code-card--image img {{
            display: block;
            margin: 0 auto;
            max-width: 100%;
            max-height: calc(480px * var(--fit-scale));
            width: auto;
            height: auto;
            object-fit: contain;
        }}
        .gen-image-card {{ margin: 0; }}
        .layout-media .slide-body {{
            display: grid;
            grid-template-columns: minmax(0, 1fr) minmax(420px, 0.92fr);
            gap: calc(42px * var(--fit-scale));
            align-items: center;
        }}
        .layout-media .media-copy,
        .layout-media .media-figure {{ min-width: 0; }}
        .gen-image-card img {{
            display: block;
            margin: 0 auto;
            max-width: 100%;
            max-height: calc(460px * var(--fit-scale));
            width: auto;
            height: auto;
            object-fit: contain;
            border-radius: 14px;
            border: 1px solid var(--border);
            box-shadow: {theme['shadow']};
        }}
        .gen-image-card figcaption {{
            margin-top: calc(12px * var(--fit-scale));
            color: var(--muted);
            font-size: calc(20px * var(--density) * var(--fit-scale));
            line-height: 1.3;
        }}
        .gen-image-labels {{
            display: flex;
            flex-wrap: wrap;
            gap: calc(7px * var(--fit-scale)) calc(12px * var(--fit-scale));
            margin: calc(10px * var(--fit-scale)) 0 0;
            padding: 0;
            list-style: none;
        }}
        .gen-image-labels li {{
            margin: 0;
            padding: calc(5px * var(--fit-scale)) calc(9px * var(--fit-scale));
            border: 1px solid var(--border);
            border-radius: 999px;
            color: var(--text-primary);
            font-size: calc(18px * var(--density) * var(--fit-scale));
        }}
        table {{
            width: 100%;
            border-collapse: collapse;
            table-layout: fixed;
        }}
        th,
        td {{
            border-bottom: 1px solid var(--border);
            padding: 12px 10px;
            text-align: left;
            vertical-align: top;
            overflow-wrap: anywhere;
        }}
        tr:last-child td {{ border-bottom: 0; }}
        pre {{
            white-space: pre-wrap;
            overflow-wrap: anywhere;
            color: var(--text-primary);
            font-family: var(--font-mono);
            font-size: calc(19px * var(--density) * var(--fit-scale));
            line-height: 1.4;
        }}
        .raw-card h3 {{
            margin-bottom: 12px;
            color: var(--accent);
            font-family: var(--font-mono);
            font-size: calc(18px * var(--fit-scale));
            letter-spacing: 0.1em;
            text-transform: uppercase;
        }}
        .raw-card pre {{ color: var(--text-secondary); }}
        .slide-progress {{
            position: fixed;
            left: 50%;
            bottom: 22px;
            transform: translateX(-50%);
            z-index: 1000;
            display: flex;
            align-items: center;
            gap: 10px;
            padding: 8px 12px;
            border-radius: 999px;
            background: rgba(0, 0, 0, 0.72);
            color: #fff;
            font: 13px/1 var(--font-mono);
            letter-spacing: 0.02em;
            opacity: 0;
            pointer-events: none;
            transition: opacity 220ms ease;
        }}
        .slide-progress.visible {{ opacity: 1; }}
        .edit-hotzone {{
            position: fixed;
            top: 0;
            left: 0;
            width: 80px;
            height: 80px;
            z-index: 10000;
        }}
        .edit-toggle {{
            position: fixed;
            top: 18px;
            left: 18px;
            z-index: 10001;
            width: 42px;
            height: 42px;
            border: 1px solid rgba(255,255,255,0.3);
            border-radius: 999px;
            background: rgba(0,0,0,0.78);
            color: #fff;
            opacity: 0;
            pointer-events: none;
            transition: opacity 180ms ease;
        }}
        .edit-toggle.show,
        .edit-toggle.active {{
            opacity: 1;
            pointer-events: auto;
        }}
        body.editing [data-editable] {{
            outline: 2px dashed var(--accent);
            outline-offset: 5px;
            cursor: text;
        }}
        .presenter-notes {{
            position: fixed;
            right: 22px;
            bottom: 22px;
            z-index: 12000;
            width: min(520px, calc(100vw - 44px));
            max-height: min(520px, calc(100vh - 44px));
            overflow: auto;
            padding: 22px 24px;
            border: 1px solid rgba(255, 255, 255, 0.24);
            border-radius: 14px;
            background: rgba(9, 12, 18, 0.94);
            box-shadow: 0 18px 60px rgba(0, 0, 0, 0.38);
            color: #f7f8fa;
            font: 17px/1.55 var(--font-body);
        }}
        .presenter-notes__label {{
            margin-bottom: 7px;
            color: #aeb8c8;
            font: 12px/1.2 var(--font-mono);
            letter-spacing: 0.14em;
            text-transform: uppercase;
        }}
        .presenter-notes h2 {{
            margin-bottom: 12px;
            color: #fff;
            font: 700 23px/1.2 var(--font-display);
        }}
        .presenter-notes__text {{ white-space: pre-wrap; }}
        @media print {{
            .presenter-notes,
            .slide-progress,
            .edit-hotzone,
            .edit-toggle {{
                visibility: hidden !important;
                opacity: 0 !important;
            }}
        }}
        html.static-export .presenter-notes {{
            display: none !important;
        }}

        /* === ENTRANCE ANIMATIONS === */
        .reveal {{
            opacity: 0;
            transform: translateY(26px);
            transition: opacity 620ms var(--ease-out-expo), transform 620ms var(--ease-out-expo);
        }}
        .slide.visible .reveal {{
            opacity: 1;
            transform: translateY(0);
        }}
        .slide.visible .reveal:nth-child(1) {{ transition-delay: 80ms; }}
        .slide.visible .reveal:nth-child(2) {{ transition-delay: 150ms; }}
        .slide.visible .reveal:nth-child(3) {{ transition-delay: 220ms; }}
        .slide.visible .reveal:nth-child(4) {{ transition-delay: 290ms; }}
        .slide.visible .reveal:nth-child(5) {{ transition-delay: 360ms; }}
        .slide.visible .reveal:nth-child(6) {{ transition-delay: 430ms; }}
        .slide.visible .reveal:nth-child(n+7) {{ transition-delay: 500ms; }}
    </style>
</head>
<body>
    <!-- Design recipe: {html.escape(design_note)} -->
    <div class="deck-viewport">
        <main class="deck-stage" id="deckStage" aria-label="{html.escape(deck.title)}">
{slides}
        </main>
    </div>
    <div class="slide-progress" id="slideProgress" aria-live="polite"></div>
    <div class="edit-hotzone" aria-hidden="true"></div>
    <button class="edit-toggle" id="editToggle" type="button" title="Edit mode (E)" aria-label="Toggle edit mode">E</button>
    <aside class="presenter-notes" id="presenterNotes" aria-live="polite" hidden>
        <p class="presenter-notes__label">Speaker notes · N to close</p>
        <h2 id="presenterNotesTitle"></h2>
        <p class="presenter-notes__text" id="presenterNotesText"></p>
    </aside>
    <script type="application/json" id="deck-manifest">{_json_for_script(manifest_for_page)}</script>
    <script type="application/json" id="speaker-notes-data">{notes_payload}</script>
    <script>
        /* === SLIDE PRESENTATION CONTROLLER === */
        class SlidePresentation {{
            constructor() {{
                this.slides = Array.from(document.querySelectorAll('.slide'));
                this.currentSlide = 0;
                this.stage = document.getElementById('deckStage');
                this.progress = document.getElementById('slideProgress');
                this.wheelLock = false;
                this.hideTimer = null;
                this.setupStageScale();
                this.setupKeyboardNav();
                this.setupTouchNav();
                this.setupWheelNav();
                this.showSlide(0);
            }}
            setupStageScale() {{
                const scale = () => {{
                    const factor = Math.min(window.innerWidth / 1920, window.innerHeight / 1080);
                    const x = (window.innerWidth - 1920 * factor) / 2;
                    const y = (window.innerHeight - 1080 * factor) / 2;
                    this.stage.style.transform = `translate(${{x}}px, ${{y}}px) scale(${{factor}})`;
                }};
                scale();
                window.addEventListener('resize', scale);
            }}
            setupKeyboardNav() {{
                document.addEventListener('keydown', (event) => {{
                    if (event.target && event.target.isContentEditable) return;
                    if (event.key === 'ArrowRight' || event.key === 'PageDown' || event.key === ' ') {{
                        event.preventDefault();
                        this.next();
                    }} else if (event.key === 'ArrowLeft' || event.key === 'PageUp') {{
                        event.preventDefault();
                        this.prev();
                    }} else if (event.key === 'Home') {{
                        event.preventDefault();
                        this.showSlide(0);
                    }} else if (event.key === 'End') {{
                        event.preventDefault();
                        this.showSlide(this.slides.length - 1);
                    }}
                }});
            }}
            setupTouchNav() {{
                let startX = 0;
                document.addEventListener('touchstart', (event) => {{
                    startX = event.changedTouches[0].clientX;
                }}, {{ passive: true }});
                document.addEventListener('touchend', (event) => {{
                    const delta = event.changedTouches[0].clientX - startX;
                    if (Math.abs(delta) > 48) delta < 0 ? this.next() : this.prev();
                }}, {{ passive: true }});
            }}
            setupWheelNav() {{
                document.addEventListener('wheel', (event) => {{
                    if (this.wheelLock || Math.abs(event.deltaY) < 42) return;
                    this.wheelLock = true;
                    event.deltaY > 0 ? this.next() : this.prev();
                    setTimeout(() => {{ this.wheelLock = false; }}, 420);
                }}, {{ passive: true }});
            }}
            showSlide(index) {{
                this.currentSlide = Math.max(0, Math.min(index, this.slides.length - 1));
                this.slides.forEach((slide, i) => {{
                    slide.classList.toggle('active', i === this.currentSlide);
                    slide.classList.toggle('visible', i === this.currentSlide);
                }});
                this.progress.textContent = `${{this.currentSlide + 1}} / ${{this.slides.length}}`;
                this.progress.classList.add('visible');
                clearTimeout(this.hideTimer);
                this.hideTimer = setTimeout(() => this.progress.classList.remove('visible'), 1400);
                const slide = this.slides[this.currentSlide];
                if (window.presenterNotes) {{
                    window.presenterNotes.showSlide(slide.dataset.slideId);
                }}
                if (window.MathJax && window.MathJax.typesetPromise) {{
                    window.MathJax.typesetPromise([slide]).then(() => fitSlide(slide)).catch(() => {{}});
                }}
            }}
            next() {{ this.showSlide(this.currentSlide + 1); }}
            prev() {{ this.showSlide(this.currentSlide - 1); }}
        }}

        /* === OFFLINE PRESENTER NOTES === */
        class PresenterNotes {{
            constructor() {{
                this.panel = document.getElementById('presenterNotes');
                this.title = document.getElementById('presenterNotesTitle');
                this.text = document.getElementById('presenterNotesText');
                this.notes = new Map();
                try {{
                    const payload = JSON.parse(
                        document.getElementById('speaker-notes-data').textContent || '[]'
                    );
                    payload.forEach((note) => this.notes.set(note.id, note));
                }} catch (error) {{}}
                this.currentSlideId = null;
                document.addEventListener('keydown', (event) => {{
                    if (
                        (event.key === 'n' || event.key === 'N') &&
                        !(event.target && event.target.isContentEditable) &&
                        !document.documentElement.classList.contains('static-export')
                    ) {{
                        event.preventDefault();
                        this.toggle();
                    }}
                }});
            }}
            showSlide(slideId) {{
                this.currentSlideId = slideId;
                const note = this.notes.get(slideId);
                this.title.textContent = note ? note.title : 'No speaker notes';
                this.text.textContent = note ? note.text : 'No correlated note is available for this slide.';
            }}
            toggle() {{
                if (!this.notes.size) return;
                this.panel.hidden = !this.panel.hidden;
            }}
        }}

        /* === INLINE EDITING CONTROLLER === */
        class InlineEditor {{
            constructor() {{
                this.isActive = false;
                this.key = `agentic-slides:${{location.pathname}}`;
                this.toggle = document.getElementById('editToggle');
                this.hotzone = document.querySelector('.edit-hotzone');
                this.hideTimeout = null;
                this.load();
                this.setup();
            }}
            setup() {{
                this.toggle.addEventListener('click', () => this.toggleEditMode());
                this.hotzone.addEventListener('mouseenter', () => {{
                    clearTimeout(this.hideTimeout);
                    this.toggle.classList.add('show');
                }});
                this.hotzone.addEventListener('mouseleave', () => {{
                    this.hideTimeout = setTimeout(() => {{
                        if (!this.isActive) this.toggle.classList.remove('show');
                    }}, 400);
                }});
                this.toggle.addEventListener('mouseenter', () => clearTimeout(this.hideTimeout));
                this.toggle.addEventListener('mouseleave', () => {{
                    this.hideTimeout = setTimeout(() => {{
                        if (!this.isActive) this.toggle.classList.remove('show');
                    }}, 400);
                }});
                document.addEventListener('keydown', (event) => {{
                    if ((event.key === 'e' || event.key === 'E') && !(event.target && event.target.isContentEditable)) {{
                        this.toggleEditMode();
                    }}
                    if ((event.ctrlKey || event.metaKey) && event.key.toLowerCase() === 's') {{
                        event.preventDefault();
                        this.save();
                    }}
                }});
                document.addEventListener('input', (event) => {{
                    if (event.target && event.target.matches('[data-editable]')) this.save();
                }});
            }}
            toggleEditMode() {{
                this.isActive = !this.isActive;
                document.body.classList.toggle('editing', this.isActive);
                this.toggle.classList.toggle('active', this.isActive);
                this.toggle.classList.add('show');
                document.querySelectorAll('[data-editable]').forEach((node) => {{
                    node.contentEditable = this.isActive ? 'true' : 'false';
                }});
            }}
            save() {{
                const values = Array.from(document.querySelectorAll('[data-editable]')).map((node) => node.innerHTML);
                localStorage.setItem(this.key, JSON.stringify(values));
            }}
            load() {{
                try {{
                    const values = JSON.parse(localStorage.getItem(this.key) || '[]');
                    document.querySelectorAll('[data-editable]').forEach((node, index) => {{
                        if (values[index]) node.innerHTML = values[index];
                    }});
                }} catch (error) {{}}
            }}
        }}

        /* === OVERFLOW AUTO-FIT ===
           Keep each Beamer frame on one HTML slide. Dense source frames may shrink
           to 60% before validation reports that the content still does not fit. */
        function fitSlide(slide) {{
            const content = slide.querySelector('.slide-content');
            const body = slide.querySelector('.slide-body');
            if (!content || !body) return;
            content.style.removeProperty('--fit-scale');
            const images = slide.querySelectorAll('.gen-image-card img');
            images.forEach((img) => img.style.removeProperty('max-height'));
            let scale = 1;
            const overflowing = () =>
                body.scrollHeight - body.clientHeight > 2 ||
                content.scrollHeight - content.clientHeight > 2;
            while (overflowing() && scale > 0.6) {{
                scale = Math.round((scale - 0.04) * 100) / 100;
                content.style.setProperty('--fit-scale', scale);
            }}
            /* Text stops at the readability floor; attached figures are
               supplementary, so they keep shrinking to absorb the rest. */
            let imageMax = 460 * scale;
            while (overflowing() && images.length && imageMax > 160) {{
                imageMax -= 20;
                images.forEach((img) => {{ img.style.maxHeight = imageMax + 'px'; }});
            }}
        }}
        function fitAllSlides() {{
            document.querySelectorAll('.slide').forEach(fitSlide);
        }}
        window.__fitAllSlides = fitAllSlides;

        window.presenterNotes = new PresenterNotes();
        window.presentation = new SlidePresentation();
        window.inlineEditor = new InlineEditor();

        window.addEventListener('load', () => {{
            const run = () => {{
                fitAllSlides();
                window.__slidesFitted = true;
            }};
            const fonts = document.fonts ? document.fonts.ready : Promise.resolve();
            fonts.then(() => {{
                if (window.MathJax && window.MathJax.startup && window.MathJax.startup.promise) {{
                    window.MathJax.startup.promise
                        .then(() => window.MathJax.typesetPromise())
                        .then(run)
                        .catch(run);
                }} else {{
                    run();
                }}
            }}).catch(run);
        }});
    </script>
</body>
</html>
"""


def render_slide(
    deck: BeamerDeck,
    slide: BeamerSlide,
    theme: dict[str, str],
    layout_rotation: list[str] | None = None,
) -> str:
    slide_id = stable_slide_id(slide.index)
    if slide.is_titlepage:
        subtitle = _preview_subtitle(deck)
        return f"""            <section class="slide title-slide active visible" id="{slide_id}" data-slide-id="{slide_id}" data-slide-index="{slide.index}">
                <div class="slide-frame">
                    <div class="slide-topline reveal"><span>{html.escape(_deck_chrome(deck))}</span><span>{deck.slide_count:02d} slides</span></div>
                    <div class="slide-content layout-hero">
                        <div class="title-lockup">
                            <h1 class="reveal" data-editable>{html.escape(deck.title)}</h1>
                            <p class="title-subtitle reveal" data-editable>{html.escape(subtitle)}</p>
                        </div>
                    </div>
                    <div class="slide-footline reveal"><span>{html.escape(_short_kicker(deck))}</span><span>{slide.index:02d} / {deck.slide_count:02d}</span></div>
                </div>
            </section>"""

    layout = choose_layout(slide, layout_rotation)
    compact = _is_compact_slide(slide)
    content_class = f"slide-content layout-{layout}" + (" compact" if compact else "")
    body = _render_body(slide, layout)
    return f"""            <section class="slide content-slide" id="{slide_id}" data-slide-id="{slide_id}" data-slide-index="{slide.index}">
                <div class="slide-frame">
                    <div class="slide-topline reveal"><span>{html.escape(_deck_chrome(deck))}</span><span>{slide.index:02d} / {deck.slide_count:02d}</span></div>
                    <div class="{content_class}">
                        {_render_head(slide)}
                        {body}
                    </div>
                    <div class="slide-footline reveal"><span>{html.escape(_short_kicker(deck))}</span><span>{html.escape(deck.title[:36])}</span></div>
                </div>
            </section>"""


def choose_layout(
    slide: BeamerSlide, layout_rotation: list[str] | None = None
) -> str:
    """Pick a per-slide layout so the deck has visual rhythm instead of one template.

    - "split": light slides keep the title-left / content-right look.
    - "math": slides with top-level equations put prose left, formulas right.
    - "top-cols": dense slides get a full-width header and a two-column body.
    - "top": everything in between - full-width header, single column below.
    """
    has_image = any(
        element.kind in {"generated_image", "user_image"}
        for element in slide.elements
    )
    has_non_image = any(
        element.kind not in {"generated_image", "user_image"}
        for element in slide.elements
    )
    if has_image and has_non_image:
        return "media"
    weight = sum(element_weight(element) for element in slide.elements)
    main, equations = _partition_equations(slide.elements)
    if equations and main and sum(element_weight(element) for element in main) <= 12:
        natural = "math"
    elif weight <= 5 and len(slide.elements) <= 2:
        natural = "split"
    elif weight > 9 or len(slide.elements) > 3 or _dominant_list_size(slide) >= 4:
        natural = "top-cols"
    else:
        natural = "top"
    if not layout_rotation:
        return natural
    allowed = [
        {"columns": "top-cols"}.get(layout, layout)
        for layout in layout_rotation
        if layout != "hero"
    ]
    if natural in allowed:
        return natural
    for fallback in ("top", "top-cols", "split", "math"):
        if fallback in allowed:
            return fallback
    return natural


def _partition_equations(elements: list[ContentElement]) -> tuple[list[ContentElement], list[ContentElement]]:
    """Hoist equations that sit at the top level or directly inside top-level blocks.

    Deeply nested equations (inside list items) stay in place so content order is preserved.
    """
    main: list[ContentElement] = []
    equations: list[ContentElement] = []
    for element in elements:
        if element.kind == "equation":
            equations.append(element)
        elif element.kind == "block" and any(child.kind == "equation" for child in element.children):
            kept = [child for child in element.children if child.kind != "equation"]
            equations.extend(child for child in element.children if child.kind == "equation")
            if kept or element.title:
                main.append(replace(element, children=kept))
        else:
            main.append(element)
    return main, equations


def _dominant_list_size(slide: BeamerSlide) -> int:
    return max((len(element.items) for element in slide.elements if element.kind == "list"), default=0)


def _render_head(slide: BeamerSlide) -> str:
    kicker, title = _split_title(slide.title or f"Slide {slide.index}")
    kicker_html = f'<p class="kicker" data-editable>{html.escape(kicker)}</p>' if kicker else ""
    return f"""<header class="slide-head reveal">
                            {kicker_html}
                            <h2 data-editable>{html.escape(title)}</h2>
                        </header>"""


def _split_title(title: str) -> tuple[str, str]:
    parts = re.split(r"\s+[-–—]\s+", title, maxsplit=1)
    if len(parts) == 2 and parts[0].strip() and parts[1].strip():
        return parts[0].strip(), parts[1].strip()
    return "", title


def _render_body(slide: BeamerSlide, layout: str) -> str:
    elements = slide.elements
    if not elements:
        return '<div class="slide-body reveal"><p data-editable></p></div>'
    if layout == "media":
        media = [
            element
            for element in elements
            if element.kind in {"generated_image", "user_image"}
        ]
        copy = [
            element
            for element in elements
            if element.kind not in {"generated_image", "user_image"}
        ]
        return (
            '<div class="slide-body">'
            f'<div class="media-copy reveal">{_render_elements(copy)}</div>'
            f'<div class="media-figure reveal">{_render_elements(media)}</div>'
            "</div>"
        )
    if layout == "math":
        main, aside = _partition_equations(elements)
        return (
            '<div class="slide-body">'
            f'<div class="col reveal">{_render_elements(main)}</div>'
            f'<div class="col col-aside reveal">{_render_elements(aside)}</div>'
            "</div>"
        )
    if layout == "top-cols":
        left, right = _split_columns(elements)
        return (
            '<div class="slide-body">'
            f'<div class="col reveal">{_render_elements(left)}</div>'
            f'<div class="col reveal">{_render_elements(right)}</div>'
            "</div>"
        )
    rendered = "\n".join(f'<div class="reveal">{render_element(element)}</div>' for element in elements)
    return f'<div class="slide-body">{rendered}</div>'


def _render_elements(elements: list[ContentElement]) -> str:
    return "\n".join(render_element(element) for element in elements)


def _split_columns(elements: list[ContentElement]) -> tuple[list[ContentElement], list[ContentElement]]:
    """Distribute elements over two columns, splitting a dominant list by items."""
    if len(elements) == 1 and elements[0].kind == "list" and len(elements[0].items) >= 2:
        return _split_list(elements[0])
    if len(elements) == 1:
        # A lone block (or similar) would leave the second column empty; break it
        # into its constituent pieces so both columns carry content.
        expanded = _expand_column_units(elements[0])
        if len(expanded) >= 2:
            elements = expanded
    else:
        # A short preamble followed by one deeply nested list is common in
        # generated slides. Expand every oversized element before balancing so
        # the preamble does not occupy one column while the full list overflows
        # the other.
        expanded = [
            unit
            for element in elements
            for unit in _expand_column_units(element)
        ]
        if len(expanded) > len(elements):
            elements = expanded
    weights = [element_weight(element) for element in elements]
    total = sum(weights)
    best_index, best_gap = 1, total
    for index in range(1, len(elements)):
        gap = abs(sum(weights[:index]) - sum(weights[index:]))
        if gap < best_gap:
            best_index, best_gap = index, gap
    left, right = elements[:best_index], elements[best_index:]
    if len(left) == 1 and left[0].kind == "list" and len(left[0].items) >= 4 and not right:
        return _split_list(left[0])
    if not right:
        return _split_columns_fallback(elements)
    return left, right


def _split_list(element: ContentElement) -> tuple[list[ContentElement], list[ContentElement]]:
    items = element.items
    weights = [item_weight(item) for item in items]
    total = sum(weights)
    mid, best_gap, running = 1, total, 0
    for index in range(1, len(items)):
        running += weights[index - 1]
        gap = abs(running - (total - running))
        if gap < best_gap:
            mid, best_gap = index, gap
    first = replace(element, items=items[:mid])
    second = replace(element, items=items[mid:], start=element.start + mid if element.ordered else 1)
    return [first], [second]


def _split_columns_fallback(elements: list[ContentElement]) -> tuple[list[ContentElement], list[ContentElement]]:
    mid = (len(elements) + 1) // 2
    return elements[:mid], elements[mid:]


_COLUMN_TARGET_WEIGHT = 10


def _expand_column_units(element: ContentElement) -> list[ContentElement]:
    """Break a lone dense element into units that can balance across two columns."""
    if element_weight(element) <= _COLUMN_TARGET_WEIGHT:
        return [element]
    if element.kind == "list":
        return _split_oversized_column_list(element)
    if element.kind == "block":
        child_units = [
            unit
            for child in element.children
            for unit in _expand_column_units(child)
        ]
        groups = _pack_column_units(child_units, _COLUMN_TARGET_WEIGHT - 1)
        if len(groups) <= 1:
            return [element]
        return [replace(element, children=group) for group in groups]
    return [element]


def _pack_column_units(
    units: list[ContentElement], cap: int
) -> list[list[ContentElement]]:
    groups: list[list[ContentElement]] = []
    current: list[ContentElement] = []
    current_weight = 0
    for unit in units:
        weight = element_weight(unit)
        if current and current_weight + weight > cap:
            groups.append(current)
            current, current_weight = [], 0
        current.append(unit)
        current_weight += weight
    if current:
        groups.append(current)
    return groups


def _split_oversized_column_list(element: ContentElement) -> list[ContentElement]:
    if len(element.items) == 1 and element.items[0].children:
        return _promote_column_list_item(element)
    chunks: list[ContentElement] = []
    items: list[ListItem] = []
    weight = 0
    start = element.start
    for item in element.items:
        weight_of_item = item_weight(item)
        if items and weight + weight_of_item > _COLUMN_TARGET_WEIGHT:
            chunks.append(replace(element, items=items, start=start))
            start = start + len(items) if element.ordered else 1
            items, weight = [], 0
        items.append(item)
        weight += weight_of_item
    if items:
        chunks.append(replace(element, items=items, start=start))
    expanded: list[ContentElement] = []
    for chunk in chunks:
        if (
            len(chunk.items) == 1
            and chunk.items[0].children
            and element_weight(chunk) > _COLUMN_TARGET_WEIGHT
        ):
            expanded.extend(_promote_column_list_item(chunk))
        else:
            expanded.append(chunk)
    return expanded


def _promote_column_list_item(element: ContentElement) -> list[ContentElement]:
    item = element.items[0]
    units: list[ContentElement] = []
    if item.text:
        units.append(replace(element, items=[replace(item, children=[])]))
    for child in item.children:
        units.extend(_expand_column_units(child))
    return units


_ALIGNMENT_SAFE_ENV_RE = re.compile(r"\\begin\{(aligned|alignedat|gathered|cases|array|[A-Za-z]*matrix)\*?\}")


_UNESCAPED_AMP_RE = re.compile(r"(?<!\\)&")


def _display_wrapper_for(env: str | None, body: str) -> str | None:
    """Pick the env needed to make an equation body valid inside \\[...\\] display math."""
    base = (env or "").rstrip("*")
    if base in {"align", "alignat", "flalign", "eqnarray"}:
        return "aligned"
    if base in {"gather", "multline"}:
        return "gathered"
    if _UNESCAPED_AMP_RE.search(body) and not _ALIGNMENT_SAFE_ENV_RE.search(body):
        # Alignment markers with no alignment environment would raise "Misplaced &".
        return "aligned"
    return None


def render_element(element: ContentElement) -> str:
    if element.kind == "text":
        return f'<p class="text-item" data-editable>{_inline_html(element.text)}</p>'
    if element.kind == "block":
        title = f"<h3 data-editable>{html.escape(element.title or 'Note')}</h3>"
        children = "\n".join(render_element(child) for child in element.children)
        if not children and element.raw:
            children = f"<p data-editable>{_inline_html(element.raw)}</p>"
        return f'<section class="content-block">{title}{children}</section>'
    if element.kind == "list":
        tag = "ol" if element.ordered else "ul"
        start = ""
        if element.ordered and element.start > 1:
            start = f' start="{element.start}" style="counter-reset: bitem {element.start - 1}"'
        return f'<{tag} class="blist"{start}>{_render_items(element.items)}</{tag}>'
    if element.kind == "equation":
        body = normalize_display_math(element.text.strip())
        wrapper = _display_wrapper_for(element.env, body)
        equation = html.escape(body)
        if wrapper:
            equation = f"\\begin{{{wrapper}}}{equation}\\end{{{wrapper}}}"
        return f'<div class="formula" data-editable>\\[{equation}\\]</div>'
    if element.kind == "table":
        return f'<div class="table-card">{_render_table(element.rows)}</div>'
    if element.kind == "code":
        if element.image_data_uri:
            label = html.escape(element.language or "code")
            return f'<figure class="code-card code-card--image"><img src="{element.image_data_uri}" alt="{label} snippet"></figure>'
        language = f" data-language=\"{html.escape(element.language)}\"" if element.language else ""
        return f'<div class="code-card"><pre{language} data-editable><code>{html.escape(element.text)}</code></pre></div>'
    if element.kind == "raw":
        label = html.escape(element.title or "LaTeX")
        return f'<div class="raw-card"><h3>{label}</h3><pre data-editable>{html.escape(element.raw)}</pre></div>'
    if element.kind in {"generated_image", "user_image"}:
        if not element.image_data_uri:
            return ""
        fallback = "Generated illustration" if element.kind == "generated_image" else "Provided image"
        title = element.title or fallback
        alt = html.escape(title)
        labels = "".join(
            f"<li>{html.escape(label)}</li>" for label in element.image_labels
        )
        legend = (
            f'<ul class="gen-image-labels" aria-label="Figure labels">{labels}</ul>'
            if labels
            else ""
        )
        caption = f"<figcaption>{html.escape(title)}</figcaption>"
        return (
            '<figure class="gen-image-card">'
            f'<img src="{element.image_data_uri}" alt="{alt}">'
            f"{caption}{legend}</figure>"
        )
    return f'<div class="raw-card"><pre data-editable>{html.escape(element.raw or element.text)}</pre></div>'


def theme_for_candidate(candidate: StyleCandidate) -> dict[str, str]:
    name_key = (candidate.slug or candidate.preset_name or candidate.name).lower()
    base = {
        "stage_bg": "#0c1016",
        "background": "#121821",
        "background_alt": "#1a2431",
        "text": "#f3efe6",
        "muted": "#aeb7c4",
        "accent": "#c8a870",
        "accent2": "#73d2de",
        "surface": "#f0ece3",
        "surface_alt": "#242f3f",
        "border": "rgba(232, 224, 210, 0.22)",
        "glow": "rgba(200, 168, 112, 0.23)",
        "grid": "rgba(255,255,255,0.05)",
        "grid_opacity": "0.55",
        "display_font": "'Source Serif 4', serif",
        "body_font": "'DM Sans', sans-serif",
        "mono_font": "'IBM Plex Mono', monospace",
        "title_size": "118px",
        "heading_size": "68px",
        "panel_border": "1px solid var(--border)",
        "panel_fill": "rgba(240, 236, 227, 0.08)",
        "shadow": "none",
        "carbon_theme": "nord",
        "carbon_background": "rgba(0,0,0,0)",
    }
    if "swiss" in name_key:
        base.update(
            background="#fbfaf7",
            background_alt="#f1f0ea",
            stage_bg="#e6e2d7",
            text="#101010",
            muted="#555b64",
            accent="#ff3300",
            accent2="#111111",
            surface="#ffffff",
            surface_alt="#ece9de",
            border="#111111",
            glow="rgba(255, 51, 0, 0.18)",
            grid="rgba(0,0,0,0.06)",
            display_font="'Archivo', sans-serif",
            body_font="'Nunito', sans-serif",
            mono_font="'IBM Plex Mono', monospace",
            panel_border="2px solid var(--border)",
            panel_fill="#ffffff",
            carbon_theme="one-light",
            carbon_background="rgba(255,255,255,1)",
        )
    elif "paper" in name_key or "cartesian" in name_key or "vellum" in name_key:
        base.update(
            background="#faf7ef",
            background_alt="#ece4d5",
            stage_bg="#161616",
            text="#1f1b17",
            muted="#6c6255",
            accent="#9e1f2f",
            accent2="#184e77",
            surface="#fffdf7",
            surface_alt="#efe7d6",
            border="#c9bda8",
            glow="rgba(158, 31, 47, 0.12)",
            grid="rgba(31,27,23,0.07)",
            display_font="'Cormorant Garamond', serif",
            body_font="'Source Serif 4', serif",
            mono_font="'IBM Plex Mono', monospace",
            panel_border="1px solid var(--border)",
            panel_fill="rgba(255, 253, 247, 0.88)",
            title_size="112px",
            heading_size="64px",
            carbon_theme="solarized light",
            carbon_background="rgba(0,0,0,0)",
        )
    elif "cobalt" in name_key:
        base.update(
            background="#f7f0df",
            background_alt="#eee2c9",
            stage_bg="#14171d",
            text="#1438ff",
            muted="#3b50a0",
            accent="#1438ff",
            accent2="#0b0d18",
            surface="#fffaf0",
            surface_alt="#e7ddc6",
            border="rgba(20, 56, 255, 0.36)",
            glow="rgba(20, 56, 255, 0.12)",
            grid="rgba(20,56,255,0.16)",
            display_font="'Playfair Display', serif",
            body_font="'IBM Plex Sans', sans-serif",
            mono_font="'IBM Plex Mono', monospace",
            panel_border="1px solid var(--border)",
            panel_fill="rgba(255, 250, 240, 0.84)",
            carbon_theme="cobalt",
            carbon_background="rgba(255,250,240,1)",
        )
    elif "block" in name_key or "neo-grid" in name_key or "raw-grid" in name_key:
        base.update(
            background="#fffdf5",
            background_alt="#f8f0d0",
            stage_bg="#111111",
            text="#000000",
            muted="#303030",
            accent="#fe90e8",
            accent2="#1438ff",
            surface="#f7cb46",
            surface_alt="#c0f7fe",
            border="#000000",
            glow="rgba(254, 144, 232, 0.2)",
            grid="rgba(0,0,0,0.08)",
            display_font="'Archivo Black', sans-serif",
            body_font="'Space Grotesk', sans-serif",
            mono_font="'Space Grotesk', monospace",
            panel_border="3px solid var(--border)",
            panel_fill="#ffffff",
            shadow="8px 8px 0 #000",
            title_size="106px",
            heading_size="62px",
            carbon_theme="one-light",
            carbon_background="rgba(255,255,255,1)",
        )
    elif "studio" in name_key or "bold signal" in name_key:
        base.update(
            background="#050505",
            background_alt="#191919",
            stage_bg="#000000",
            text="#f4f1e8",
            muted="#a9a9a0",
            accent="#eaff00",
            accent2="#ff5722",
            surface="#eaff00",
            surface_alt="#202020",
            border="rgba(234, 255, 0, 0.32)",
            glow="rgba(234, 255, 0, 0.2)",
            grid="rgba(234,255,0,0.08)",
            display_font="'Archivo Black', sans-serif",
            body_font="'Space Grotesk', sans-serif",
            mono_font="'IBM Plex Mono', monospace",
            carbon_theme="seti",
            carbon_background="rgba(0,0,0,0)",
        )
    elif "chalkboard" in name_key or "vector" in name_key:
        base.update(
            background="#14231b",
            background_alt="#0d1812",
            stage_bg="#080f0b",
            text="#f2f5ee",
            muted="#a8bcab",
            accent="#ffd66b",
            accent2="#8fd0ff",
            surface="#1b2f24",
            surface_alt="#22392c",
            border="rgba(242, 245, 238, 0.28)",
            glow="rgba(255, 214, 107, 0.16)",
            grid="rgba(242,245,238,0.07)",
            display_font="'Space Grotesk', sans-serif",
            body_font="'Nunito', sans-serif",
            mono_font="'IBM Plex Mono', monospace",
            panel_border="1px dashed var(--border)",
            panel_fill="rgba(242, 245, 238, 0.06)",
            carbon_theme="nord",
            carbon_background="rgba(0,0,0,0)",
        )
    return base


def _font_links(theme: dict[str, str]) -> str:
    return (
        '<link rel="preconnect" href="https://fonts.googleapis.com">\n'
        '    <link rel="preconnect" href="https://fonts.gstatic.com" crossorigin>\n'
        '    <link href="https://fonts.googleapis.com/css2?'
        'family=Archivo:wght@700;800;900&family=Archivo+Black&family=Cormorant+Garamond:wght@400;600;700&'
        'family=DM+Sans:wght@400;500;700&family=IBM+Plex+Mono:wght@400;500;700&family=IBM+Plex+Sans:wght@400;500;700&'
        'family=Nunito:wght@400;500;700&family=Playfair+Display:wght@600;700&family=Source+Serif+4:wght@400;600;700&'
        'family=Space+Grotesk:wght@400;500;700&display=swap" rel="stylesheet">'
    )


def _render_items(items: list[ListItem]) -> str:
    rendered = []
    for item in items:
        children = "\n".join(render_element(child) for child in item.children)
        rendered.append(f"<li data-editable>{_inline_html(item.text)}{children}</li>")
    return "\n".join(rendered)


def _render_table(rows: list[list[str]]) -> str:
    if not rows:
        return "<pre data-editable></pre>"
    body = []
    for row in rows:
        cells = "".join(f"<td data-editable>{_inline_html(cell)}</td>" for cell in row)
        body.append(f"<tr>{cells}</tr>")
    return f"<table>{''.join(body)}</table>"


_INLINE_MATH_RE = re.compile(r"\\\((.+?)\\\)|\\\[(.+?)\\\]", re.DOTALL)


_INLINE_MATH_TOKEN_RE = re.compile("(\\d+)")


def _inline_html(text: str) -> str:
    # Protect MathJax spans so escaping and the markdown-ish regexes below
    # cannot mangle TeX such as `x^*` or `a * b`.
    spans: list[str] = []

    def stash(match: re.Match[str]) -> str:
        spans.append(match.group(0))
        return f"{len(spans) - 1}"

    text = _INLINE_MATH_RE.sub(stash, text)
    escaped = html.escape(text)
    escaped = re.sub(r"\*\*(.+?)\*\*", r"<strong>\1</strong>", escaped)
    escaped = re.sub(r"\*(.+?)\*", r"<em>\1</em>", escaped)
    escaped = escaped.replace("\n", "<br>")
    return _INLINE_MATH_TOKEN_RE.sub(lambda match: html.escape(spans[int(match.group(1))]), escaped)


def _is_compact_slide(slide: BeamerSlide) -> bool:
    total = sum(element_weight(element) for element in slide.elements)
    return total > 9 or len(slide.elements) > 3


def _preview_subtitle(deck: BeamerDeck) -> str:
    for slide in deck.slides:
        if slide.is_titlepage:
            continue
        if slide.title and slide.title != deck.title:
            return slide.title
    return "A faithful interactive HTML version of the source Beamer deck"


def _short_kicker(deck: BeamerDeck) -> str:
    stem = deck.source_path.stem.replace("-", " ").title()
    return stem[:32]


def _deck_chrome(deck: BeamerDeck) -> str:
    return deck.metadata.get("author") or deck.source_path.stem.replace("-", " ").title()




def _json_for_script(value: object) -> str:
    """Serialize JSON without allowing data to terminate its script element."""
    return (
        json.dumps(value, ensure_ascii=False, separators=(",", ":"))
        .replace("&", "\\u0026")
        .replace("<", "\\u003c")
        .replace(">", "\\u003e")
        .replace("\u2028", "\\u2028")
        .replace("\u2029", "\\u2029")
    )


def _indent(text: str, spaces: int) -> str:
    prefix = " " * spaces
    return "\n".join(prefix + line if line else "" for line in text.splitlines())


# ---------------------------------------------------------------------------
# Offline runtime
# ---------------------------------------------------------------------------

def runtime_asset_root() -> Path:
    return slide_gen_asset_root() / "runtime"


def prepare_offline_runtime(
    chapter_dir: Path, style: CourseSlideStyle
) -> tuple[Path, str]:
    """Copy only the selected runtime assets and return their @font-face CSS."""
    source_root = runtime_asset_root()
    output_root = chapter_dir / "html" / "assets"
    if output_root.exists():
        shutil.rmtree(output_root)
    math_output = output_root / "mathjax"
    font_output = output_root / "fonts"
    license_output = output_root / "licenses"
    math_output.mkdir(parents=True, exist_ok=True)
    font_output.mkdir(parents=True, exist_ok=True)
    license_output.mkdir(parents=True, exist_ok=True)

    math_source = source_root / "mathjax" / "tex-svg.js"
    if not math_source.is_file():
        raise FrontendSlidesError(f"Offline MathJax runtime is missing: {math_source}")
    shutil.copy2(math_source, math_output / math_source.name)
    math_license = source_root / "licenses" / "MATHJAX-LICENSE.txt"
    if math_license.is_file():
        shutil.copy2(math_license, license_output / math_license.name)

    selected_fonts = {
        style.render_theme.display_font,
        style.render_theme.body_font,
        style.render_theme.mono_font,
    }
    css_blocks: list[str] = []
    for font_id in sorted(selected_fonts):
        try:
            family, stem = FONT_FAMILIES[font_id]
        except KeyError as exc:
            raise FrontendSlidesError(f"Unsupported packaged font: {font_id}") from exc
        for weight in (400, 700):
            filename = f"{stem}-{weight}.woff2"
            source = source_root / "fonts" / filename
            if not source.is_file():
                raise FrontendSlidesError(f"Packaged font file is missing: {source}")
            shutil.copy2(source, font_output / filename)
            css_blocks.append(
                "@font-face {"
                f"font-family:'{family}';"
                f"src:url('assets/fonts/{filename}') format('woff2');"
                f"font-style:normal;font-weight:{weight};font-display:swap;"
                "}"
            )

        license_key = {
            "archivo": "OFL-ARCHIVO.txt",
            "dm-sans": "OFL-DM-SANS.txt",
            "source-serif-4": "OFL-SOURCE-SERIF-4.txt",
            "space-grotesk": "OFL-SPACE-GROTESK.txt",
            "ibm-plex-mono": "OFL-IBM-PLEX-MONO.txt",
        }[font_id]
        license_source = source_root / "licenses" / license_key
        if license_source.is_file():
            shutil.copy2(license_source, license_output / license_source.name)

    return output_root, "<style>\n" + "\n".join(css_blocks) + "\n</style>"


def offline_runtime_complete(chapter_dir: Path, style: CourseSlideStyle) -> bool:
    """Return whether every runtime file referenced by this course style exists."""
    output_root = chapter_dir / "html" / "assets"
    required = [output_root / "mathjax" / "tex-svg.js"]
    selected_fonts = {
        style.render_theme.display_font,
        style.render_theme.body_font,
        style.render_theme.mono_font,
    }
    for font_id in selected_fonts:
        try:
            _, stem = FONT_FAMILIES[font_id]
        except KeyError:
            return False
        required.extend(
            output_root / "fonts" / f"{stem}-{weight}.woff2"
            for weight in (400, 700)
        )
    return all(path.is_file() and path.stat().st_size > 0 for path in required)


# ---------------------------------------------------------------------------
# Validation
# ---------------------------------------------------------------------------

def validate_html_contract(html: str, expected_slide_count: int, viewport_css: str) -> list[str]:
    errors: list[str] = []
    slide_count = len(re.findall(r"<section\s+class=\"[^\"]*\bslide\b", html))
    if slide_count != expected_slide_count:
        errors.append(f"Expected {expected_slide_count} .slide sections, found {slide_count}.")
    if "tex-svg.js" not in html or "MathJax" not in html:
        errors.append("MathJax configuration is missing.")
    viewport_probe = "FIXED 16:9 STAGE: MANDATORY BASE STYLES"
    if viewport_probe not in html or viewport_probe not in viewport_css:
        errors.append("Full viewport-base.css contents do not appear to be included.")
    if re.search(r"\.slide[^{]*\{[^}]*display\s*:\s*none", html, flags=re.IGNORECASE | re.DOTALL):
        errors.append("Slide switching must not use display:none.")
    if "width: 1920px" not in html or "height: 1080px" not in html:
        errors.append("Fixed 1920x1080 stage dimensions are missing.")
    errors.extend(validate_math_syntax(html))
    return errors


def validate_offline_contract(html: str) -> list[str]:
    errors: list[str] = []
    if re.search(r"(?:src|href)\s*=\s*[\"']https?://", html, flags=re.IGNORECASE):
        errors.append("Offline slides must not load scripts, styles, fonts, or images over HTTP.")
    required = (
        "assets/mathjax/tex-svg.js",
        "assets/fonts/",
    )
    for reference in required:
        if reference not in html:
            errors.append(f"Offline runtime reference is missing: {reference}")
    return errors


_MARKUP_BLOCK_RE = re.compile(r"<script.*?</script>|<style.*?</style>", re.DOTALL | re.IGNORECASE)


_MATH_SPAN_RE = re.compile(r"\\\[(.*?)\\\]|\\\((.*?)\\\)", re.DOTALL)


_BAD_NESTED_ENV_RE = re.compile(r"\\begin\{(align|alignat|flalign|eqnarray|gather|multline|equation)\*?\}")


_ALIGNMENT_HOST_ENV_RE = re.compile(r"\\begin\{(aligned|alignedat|gathered|cases|array|[A-Za-z]*matrix)\*?\}")


_UNESCAPED_AMP_ENTITY_RE = re.compile(r"(?<!\\)&amp;")


def validate_math_syntax(html: str) -> list[str]:
    """Statically reject TeX that MathJax is guaranteed to render as an error box."""
    errors: list[str] = []
    visible = _MARKUP_BLOCK_RE.sub("", html)
    for match in _MATH_SPAN_RE.finditer(visible):
        body = match.group(1) or match.group(2) or ""
        excerpt = " ".join(body.split())[:90]
        bad_env = _BAD_NESTED_ENV_RE.search(body)
        if bad_env:
            errors.append(
                f"Math span nests \\begin{{{bad_env.group(1)}}}, which MathJax rejects"
                f" inside \\[...\\]: {excerpt!r}"
            )
        elif _UNESCAPED_AMP_ENTITY_RE.search(body) and not _ALIGNMENT_HOST_ENV_RE.search(body):
            errors.append(
                f"Math span uses '&' alignment without an alignment environment"
                f" (MathJax 'Misplaced &'): {excerpt!r}"
            )
    return errors




def validate_with_playwright(html_path: Path, expected_slide_count: int) -> list[str]:
    """Best-effort visual validation. Returns a skip-style warning if browsers are unavailable."""
    try:
        from playwright.sync_api import Error as PlaywrightError
        from playwright.sync_api import sync_playwright
    except Exception as exc:
        return [f"Playwright import unavailable: {exc}"]

    try:
        with sync_playwright() as p:
            browser = p.chromium.launch()
            page = browser.new_page(viewport={"width": 1280, "height": 720})
            page.goto(html_path.resolve().as_uri())
            page.wait_for_load_state("load")
            try:
                page.wait_for_function("window.__slidesFitted === true", timeout=20000)
            except Exception:
                pass  # older decks without the auto-fit hook; probe re-runs fit itself
            slide_count = page.locator(".slide").count()
            visible_count = page.locator(".slide.visible").count()
            stage_box = page.locator(".deck-stage").bounding_box()
            overflowing = page.evaluate(_OVERFLOW_PROBE_JS)
            math_errors = page.evaluate(_MATHJAX_ERROR_PROBE_JS)
            browser.close()
    except PlaywrightError as exc:
        return [f"Playwright browser unavailable: {exc}"]
    except Exception as exc:
        return [f"Playwright validation failed: {exc}"]

    errors = []
    if slide_count != expected_slide_count:
        errors.append(f"Expected {expected_slide_count} slides in browser, found {slide_count}.")
    if visible_count != 1:
        errors.append(f"Expected one visible slide, found {visible_count}.")
    if not stage_box or stage_box["width"] <= 0 or stage_box["height"] <= 0:
        errors.append("Deck stage did not render with positive dimensions.")
    if overflowing:
        errors.append(
            "Slide content overflows the fixed 1920x1080 stage on slides: "
            + ", ".join(str(index) for index in overflowing)
        )
    for math_error in math_errors:
        errors.append(f"MathJax error on slide {math_error['slide']}: {math_error['message']}")
    return errors


_MATHJAX_ERROR_PROBE_JS = """
() => {
    const seen = [];
    document.querySelectorAll('mjx-merror, [data-mjx-error]').forEach((node) => {
        const slide = node.closest('.slide');
        seen.push({
            slide: slide ? Number(slide.dataset.slideIndex ?? -1) : -1,
            message: node.getAttribute('data-mjx-error') || node.textContent.trim(),
        });
    });
    return seen;
}
"""


_OVERFLOW_PROBE_JS = """
() => {
    if (window.__fitAllSlides) window.__fitAllSlides();
    const overflowing = [];
    document.querySelectorAll('.slide').forEach((slide, index) => {
        const regions = [
            slide.querySelector('.slide-body'),
            slide.querySelector('.slide-content'),
            slide,
        ].filter(Boolean);
        const spills = regions.some(
            (region) =>
                region.scrollHeight - region.clientHeight > 2 ||
                region.scrollWidth - region.clientWidth > 2
        );
        if (spills) overflowing.push(index + 1);
    });
    return overflowing;
}
"""


# ---------------------------------------------------------------------------
# PDF and PPTX export
# ---------------------------------------------------------------------------

WIDTH = 1920


HEIGHT = 1080


def export_html_deck(
    html_path: Path,
    *,
    pdf_path: Path | None,
    pptx_path: Path | None,
) -> None:
    """Capture the deck once and create whichever static exports are requested."""
    if pdf_path is None and pptx_path is None:
        return
    if not html_path.is_file():
        raise FrontendSlidesError(f"Presentation HTML not found: {html_path}")
    with tempfile.TemporaryDirectory(prefix="instructional-slides-export-") as tmp:
        screenshot_dir = Path(tmp) / "screenshots"
        screenshot_dir.mkdir()
        screenshots = capture_slide_screenshots(html_path, screenshot_dir)
        if pdf_path is not None:
            _write_pdf(screenshots, pdf_path)
        if pptx_path is not None:
            _write_pptx(screenshots, pptx_path)


def capture_slide_screenshots(html_path: Path, screenshot_dir: Path) -> list[Path]:
    try:
        from playwright.sync_api import Error as PlaywrightError
        from playwright.sync_api import sync_playwright
    except Exception as exc:
        raise FrontendSlidesError(f"Playwright is unavailable for export: {exc}") from exc

    screenshots: list[Path] = []
    try:
        with sync_playwright() as playwright:
            browser = playwright.chromium.launch()
            page = browser.new_page(viewport={"width": WIDTH, "height": HEIGHT})
            page.goto(html_path.resolve().as_uri(), wait_until="load")
            try:
                page.wait_for_function("window.__slidesFitted === true", timeout=30000)
            except Exception:
                pass
            page.add_style_tag(content=_EXPORT_MODE_CSS)
            page.evaluate("document.documentElement.classList.add('static-export')")
            slide_count = page.locator(".slide").count()
            if slide_count < 1:
                raise FrontendSlidesError("No .slide elements were found in slides.html.")
            stage = page.locator(".deck-stage")
            if stage.count() != 1:
                raise FrontendSlidesError("Expected exactly one .deck-stage in slides.html.")
            for index in range(slide_count):
                page.evaluate(_SHOW_SLIDE_JS, index)
                page.evaluate(_WAIT_FOR_PAINT_JS)
                output = screenshot_dir / f"slide-{index + 1:03d}.png"
                stage.screenshot(path=str(output))
                screenshots.append(output)
            browser.close()
    except FrontendSlidesError:
        raise
    except PlaywrightError as exc:
        raise FrontendSlidesError(f"Playwright Chromium export failed: {exc}") from exc
    except Exception as exc:
        raise FrontendSlidesError(f"HTML slide capture failed: {exc}") from exc
    return screenshots


def _write_pdf(images: list[Path], output: Path) -> None:
    try:
        from playwright.sync_api import Error as PlaywrightError
        from playwright.sync_api import sync_playwright
    except Exception as exc:
        raise FrontendSlidesError(f"Playwright is unavailable for PDF export: {exc}") from exc

    pages = []
    for image in images:
        encoded = base64.b64encode(image.read_bytes()).decode("ascii")
        pages.append(
            '<section class="page">'
            f'<img src="data:image/png;base64,{encoded}" alt="">'
            "</section>"
        )
    document = f"""<!doctype html>
<html><head><meta charset="utf-8"><style>
* {{ box-sizing:border-box; margin:0; padding:0; }}
@page {{ size:{WIDTH}px {HEIGHT}px; margin:0; }}
.page {{ width:{WIDTH}px; height:{HEIGHT}px; page-break-after:always; overflow:hidden; }}
.page:last-child {{ page-break-after:auto; }}
img {{ display:block; width:100%; height:100%; object-fit:contain; }}
</style></head><body>{''.join(pages)}</body></html>"""
    output.parent.mkdir(parents=True, exist_ok=True)
    try:
        with sync_playwright() as playwright:
            browser = playwright.chromium.launch()
            page = browser.new_page()
            page.set_content(document, wait_until="load")
            page.pdf(
                path=str(output),
                width=f"{WIDTH}px",
                height=f"{HEIGHT}px",
                print_background=True,
                margin={"top": "0", "right": "0", "bottom": "0", "left": "0"},
            )
            browser.close()
    except PlaywrightError as exc:
        raise FrontendSlidesError(f"Playwright PDF export failed: {exc}") from exc
    except Exception as exc:
        raise FrontendSlidesError(f"HTML PDF export failed: {exc}") from exc


def _write_pptx(images: list[Path], output: Path) -> None:
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except Exception as exc:
        raise FrontendSlidesError(f"python-pptx is unavailable: {exc}") from exc

    try:
        presentation = Presentation()
        presentation.slide_width = Inches(13.333333)
        presentation.slide_height = Inches(7.5)
        while presentation.slides:
            relation_id = presentation.slides._sldIdLst[0].rId
            presentation.part.drop_rel(relation_id)
            del presentation.slides._sldIdLst[0]
        layout = presentation.slide_layouts[6]
        for image in images:
            slide = presentation.slides.add_slide(layout)
            slide.shapes.add_picture(
                str(image),
                0,
                0,
                width=presentation.slide_width,
                height=presentation.slide_height,
            )
        output.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(output)
    except Exception as exc:
        raise FrontendSlidesError(f"Static PPTX export failed: {exc}") from exc


_SHOW_SLIDE_JS = """
(index) => {
    const slides = Array.from(document.querySelectorAll('.slide'));
    if (window.presentation && typeof window.presentation.showSlide === 'function') {
        window.presentation.showSlide(index);
    }
    slides.forEach((slide, current) => {
        slide.classList.toggle('active', current === index);
        slide.classList.toggle('visible', current === index);
        slide.style.opacity = current === index ? '1' : '0';
        slide.style.visibility = current === index ? 'visible' : 'hidden';
        slide.style.pointerEvents = current === index ? 'auto' : 'none';
    });
    if (window.__fitAllSlides) window.__fitAllSlides();
    const selected = slides[index];
    if (selected) {
        selected.querySelectorAll('.reveal').forEach((node) => {
            node.style.opacity = '1';
            node.style.transform = 'none';
            node.style.visibility = 'visible';
        });
    }
}
"""


_EXPORT_MODE_CSS = """
html.static-export *,
html.static-export *::before,
html.static-export *::after {
    animation: none !important;
    transition: none !important;
    transition-delay: 0s !important;
}
html.static-export .reveal {
    opacity: 1 !important;
    transform: none !important;
    visibility: visible !important;
}
html.static-export .slide-progress,
html.static-export .deck-controls,
html.static-export .edit-hotzone,
html.static-export .edit-toggle {
    display: none !important;
}
"""


_WAIT_FOR_PAINT_JS = """
() => new Promise((resolve) => {
    requestAnimationFrame(() => requestAnimationFrame(resolve));
})
"""


# ---------------------------------------------------------------------------
# Chapter finalization
# ---------------------------------------------------------------------------

MANIFEST_FILENAME = "frontend-slides-manifest.json"


MANIFEST_SCHEMA_VERSION = 3


LEGACY_SPLIT_REPORT_FILENAME = "slide-splits.json"


def finalize_chapter(
    course_dir: Path | str,
    chapter_dir: Path | str,
    *,
    llm: Any | None = None,
    chapter: dict[str, str] | None = None,
    image_config: ImageGenerationConfig | None = None,
) -> ChapterFrontendResult:
    """Compile LaTeX and deterministically create resumable offline frontend artifacts."""
    course_path = Path(course_dir)
    chapter_path = Path(chapter_dir)
    tex_path = chapter_path / "slides.tex"
    script_path = chapter_path / "script.md"
    latex_pdf_path = chapter_path / "slides.pdf"
    html_path = chapter_path / "html" / "slides.html"
    html_pdf_path = chapter_path / "slides-html.pdf"
    html_pptx_path = chapter_path / "slides-html.pptx"
    manifest_path = chapter_path / MANIFEST_FILENAME
    legacy_split_report_path = chapter_path / LEGACY_SPLIT_REPORT_FILENAME

    if not tex_path.is_file() or tex_path.stat().st_size == 0:
        raise FrontendSlidesError(f"Chapter LaTeX source is missing or empty: {tex_path}")
    if not script_path.is_file() or script_path.stat().st_size == 0:
        raise FrontendSlidesError(f"Chapter speaker notes are missing or empty: {script_path}")
    preflight = normalize_beamer_file(tex_path)
    if preflight.removed_list_wrapper_pairs:
        print(
            "[preflight] Repaired slides.tex by flattening "
            f"{preflight.removed_list_wrapper_pairs} list wrapper(s) beyond "
            "Beamer's 3-level nesting limit."
        )
    if preflight.inserted_list_closures:
        print(
            "[preflight] Repaired slides.tex by closing "
            f"{preflight.inserted_list_closures} unclosed generated list "
            "environment(s)."
        )
    if preflight.injected_color_definitions:
        print(
            "[preflight] Repaired slides.tex by auto-defining missing color(s): "
            f"{', '.join(preflight.injected_color_definitions)}"
        )
    if preflight.repaired_nested_math_environments:
        print(
            "[preflight] Repaired slides.tex by normalizing "
            f"{preflight.repaired_nested_math_environments} nested display-math "
            "environment(s)."
        )
    style = load_course_slide_style(course_path)
    image_config_warnings: list[str] = []
    if image_config is None:
        try:
            image_config = load_image_generation_config(course_path)
        except ValueError as exc:
            image_config = ImageGenerationConfig()
            image_config_warnings.append(
                f"Invalid image-generation configuration was treated as disabled: {exc}"
            )
    image_config = image_config.validated()
    style_path = course_path / STYLE_FILENAME
    source_hash = sha256_file(tex_path)
    initial_script_hash = sha256_file(script_path)
    style_hash = sha256_file(style_path)
    previous = _read_manifest(manifest_path)
    manifest_matches = previous.get("schema_version") == MANIFEST_SCHEMA_VERSION
    source_matches = previous.get("source_sha256") == source_hash
    script_matches = previous.get("script_sha256") == initial_script_hash
    style_matches = previous.get("style_sha256") == style_hash
    inputs_match = (
        manifest_matches and source_matches and script_matches and style_matches
    )
    image_fingerprint = image_request_fingerprint(
        source_sha256=source_hash,
        style_sha256=style_hash,
        chapter=chapter,
        config=image_config,
        guidance=style.image_guidance,
    )
    wants_images = desired_image_state(image_config, style.image_guidance)
    raw_previous_image_state = previous.get("images")
    previous_image_state = (
        raw_previous_image_state
        if isinstance(raw_previous_image_state, dict)
        else {}
    )
    previous_generated = previous_image_state.get("generated", 0)
    if not wants_images and raw_previous_image_state is None:
        # Legacy manifests are already in the desired default-disabled state.
        image_state_current = True
    else:
        image_state_current = (
            previous_image_state.get("request_fingerprint") == image_fingerprint
            and isinstance(previous_generated, int)
            and (
                image_cache_is_current(chapter_path, image_fingerprint)
                if wants_images
                else previous_generated == 0
            )
        )
    if image_config.replace_images:
        image_state_current = False
    legacy_split_report_path.unlink(missing_ok=True)

    if not source_matches or not _nonempty(latex_pdf_path):
        try:
            LaTeXCompiler(str(chapter_path)).compile_one(tex_path)
        except Exception as exc:
            raise FrontendSlidesError(f"LaTeX PDF compilation failed: {exc}") from exc
    if not _nonempty(latex_pdf_path):
        raise FrontendSlidesError(f"LaTeX compilation did not produce {latex_pdf_path}.")

    runtime_ok = offline_runtime_complete(chapter_path, style)
    complete = all(
        _nonempty(path)
        for path in (
            html_path,
            html_pdf_path,
            html_pptx_path,
            manifest_path,
        )
    )
    if inputs_match and runtime_ok and complete and image_state_current:
        cached_run = ChapterImageResult(
            generated=int(previous_generated),
            reused_from_cache=bool(wants_images),
            committed=True,
            slides=[
                int(index)
                for index in previous_image_state.get("slide_indexes", [])
                if isinstance(index, int) and not isinstance(index, bool)
            ],
            request_fingerprint=image_fingerprint,
        )
        try:
            append_image_statistics(chapter_path, cached_run)
        except OSError as exc:
            print(f"Warning: could not record image-generation statistics: {exc}")
        _remove_legacy_frontend_layout(chapter_path)
        return ChapterFrontendResult(
            html_path=html_path,
            latex_pdf_path=latex_pdf_path,
            html_pdf_path=html_pdf_path,
            html_pptx_path=html_pptx_path,
            manifest_path=manifest_path,
            slide_count=int(previous.get("slide_count", 0)),
            skipped=True,
        )

    assets = load_assets()
    deck = parse_beamer(tex_path)
    script_markdown = script_path.read_text(encoding="utf-8")
    try:
        canonical_script = upgrade_legacy_script(deck, script_markdown)
    except FrontendSlidesError as exc:
        raise FrontendSlidesError(
            f"{exc} Run the chapter note-only repair path before finalizing this deck."
        ) from exc
    if canonical_script != script_markdown:
        _atomic_write(script_path, canonical_script)
        inputs_match = False
    script_hash = sha256_file(script_path)
    speaker_notes = correlate_speaker_notes(
        deck, script_path.read_text(encoding="utf-8")
    )
    image_result = augment_deck_with_generated_images(
        deck,
        chapter_path=chapter_path,
        style=style,
        chapter=chapter,
        llm=llm,
        config=image_config,
        source_sha256=source_hash,
        style_sha256=style_hash,
    )
    image_result.warnings[:0] = image_config_warnings

    html_stale = (
        not inputs_match
        or not _nonempty(html_path)
        or not runtime_ok
        or not image_state_current
        or image_config.replace_images
    )
    staged_html_path: Path | None = None
    if html_stale:
        html_temporary = html_path.with_name("slides.tmp.html")
        try:
            _, font_css = prepare_offline_runtime(chapter_path, style)
            html = render_course_presentation_html(
                deck,
                style,
                assets,
                font_css=font_css,
                speaker_notes=speaker_notes,
            )
            errors = validate_html_contract(
                html, deck.slide_count, assets.viewport_css
            )
            errors.extend(validate_offline_contract(html))
            if errors:
                raise FrontendSlidesError("; ".join(errors))
            html_temporary.write_text(html, encoding="utf-8")
            visual_errors = validate_with_playwright(
                html_temporary, deck.slide_count
            )
            if visual_errors:
                raise FrontendSlidesError("; ".join(visual_errors))
            staged_html_path = html_temporary
        except Exception:
            discard_image_result(image_result)
            html_temporary.unlink(missing_ok=True)
            image_result.warnings.append(
                "Frontend validation failed; staged image changes were discarded."
            )
            try:
                append_image_statistics(chapter_path, image_result)
            except OSError:
                pass
            raise

    need_pdf = html_stale or not _nonempty(html_pdf_path)
    need_pptx = html_stale or not _nonempty(html_pptx_path)
    pdf_temporary = (
        chapter_path / "slides-html.tmp.pdf" if need_pdf else None
    )
    pptx_temporary = (
        chapter_path / "slides-html.tmp.pptx" if need_pptx else None
    )
    try:
        export_html_deck(
            staged_html_path or html_path,
            pdf_path=pdf_temporary,
            pptx_path=pptx_temporary,
        )
        if pdf_temporary is not None:
            if not _nonempty(pdf_temporary):
                raise FrontendSlidesError("HTML PDF exporter produced no output.")
        if pptx_temporary is not None:
            if not _nonempty(pptx_temporary):
                raise FrontendSlidesError("HTML PPTX exporter produced no output.")
        if staged_html_path is not None:
            staged_html_path.replace(html_path)
        if pdf_temporary is not None:
            pdf_temporary.replace(html_pdf_path)
        if pptx_temporary is not None:
            pptx_temporary.replace(html_pptx_path)
    except Exception as exc:
        image_transaction = (
            image_config.replace_images
            or image_result.pending_manifest is not None
        )
        discard_image_result(image_result)
        image_result.warnings.append(
            "Frontend export failed; staged image changes were discarded."
        )
        if image_transaction:
            if staged_html_path is not None:
                staged_html_path.unlink(missing_ok=True)
            if pdf_temporary is not None:
                pdf_temporary.unlink(missing_ok=True)
            if pptx_temporary is not None:
                pptx_temporary.unlink(missing_ok=True)
            message = "prior artifacts were preserved"
        else:
            if staged_html_path is not None and staged_html_path.is_file():
                staged_html_path.replace(html_path)
            if pdf_temporary is not None and _valid_pdf(pdf_temporary):
                pdf_temporary.replace(html_pdf_path)
            elif pdf_temporary is not None:
                pdf_temporary.unlink(missing_ok=True)
            if pptx_temporary is not None and _valid_pptx(pptx_temporary):
                pptx_temporary.replace(html_pptx_path)
            elif pptx_temporary is not None:
                pptx_temporary.unlink(missing_ok=True)
            message = "successful artifacts were preserved"
        try:
            append_image_statistics(chapter_path, image_result)
        except OSError:
            pass
        raise FrontendSlidesError(
            f"Frontend export failed; {message}: {exc}"
        ) from exc
    try:
        commit_image_result(chapter_path, image_result)
    except OSError as exc:
        discard_image_result(image_result)
        image_result.committed = False
        image_result.warnings.append(
            f"Image manifest commit failed; the next run will retry: {exc}"
        )

    warnings = []
    if preflight.removed_list_wrapper_pairs:
        warnings.append(
            "LaTeX preflight flattened "
            f"{preflight.removed_list_wrapper_pairs} list wrapper(s) beyond "
            "Beamer's 3-level nesting limit."
        )
    if preflight.inserted_list_closures:
        warnings.append(
            "LaTeX preflight closed "
            f"{preflight.inserted_list_closures} unclosed generated list "
            "environment(s)."
        )
    if preflight.injected_color_definitions:
        warnings.append(
            "LaTeX preflight auto-defined missing color(s): "
            + ", ".join(preflight.injected_color_definitions)
        )
    if preflight.repaired_nested_math_environments:
        warnings.append(
            "LaTeX preflight normalized "
            f"{preflight.repaired_nested_math_environments} nested display-math "
            "environment(s)."
        )
    if deck.unsupported_environments:
        warnings.append(
            "Unsupported LaTeX environments were preserved as source cards: "
            + ", ".join(deck.unsupported_environments)
        )
    warnings.extend(image_result.warnings)
    current_images = [record.manifest_dict() for record in image_result.images]
    manifest = {
        "schema_version": MANIFEST_SCHEMA_VERSION,
        "source": tex_path.name,
        "source_sha256": source_hash,
        "script": script_path.name,
        "script_sha256": script_hash,
        "style_sha256": style_hash,
        "selected_style": asdict(style.selected_style),
        "slide_count": deck.slide_count,
        "offline": True,
        "runtime": "html/assets",
        "warnings": warnings,
        "speaker_notes": notes_manifest(speaker_notes),
        "images": {
            "pipeline_version": IMAGE_PIPELINE_VERSION,
            "enabled": image_config.enabled,
            "guidance_enabled": style.image_guidance.enabled,
            "effective": wants_images,
            "effective_cap": min(
                image_config.effective_operator_cap,
                style.image_guidance.max_images_per_chapter,
            ),
            "model": image_config.model,
            "size": image_config.size,
            "quality": image_config.quality,
            "request_fingerprint": image_fingerprint,
            "generated": len(current_images),
            "slide_indexes": image_result.slides,
            "records": current_images,
            "reused_from_cache": image_result.reused_from_cache,
            "replacement_requested": image_config.replace_images,
            "committed": image_result.committed,
            "incremental_estimated_cost_usd": image_result.estimated_cost_usd,
            "pipeline_warnings": image_result.warnings,
        },
        "artifacts": {
            "latex_pdf": latex_pdf_path.name,
            "html": html_path.relative_to(chapter_path).as_posix(),
            "html_pdf": html_pdf_path.name,
            "html_pptx": html_pptx_path.name,
        },
        "artifact_sha256": {
            "latex_pdf": sha256_file(latex_pdf_path),
            "html": sha256_file(html_path),
            "html_pdf": sha256_file(html_pdf_path),
            "html_pptx": sha256_file(html_pptx_path),
        },
    }
    _atomic_write(
        manifest_path, json.dumps(manifest, indent=2, sort_keys=True) + "\n"
    )
    try:
        append_image_statistics(chapter_path, image_result)
    except OSError as exc:
        print(f"Warning: could not record image-generation statistics: {exc}")
    _remove_legacy_frontend_layout(chapter_path)
    return ChapterFrontendResult(
        html_path=html_path,
        latex_pdf_path=latex_pdf_path,
        html_pdf_path=html_pdf_path,
        html_pptx_path=html_pptx_path,
        manifest_path=manifest_path,
        slide_count=deck.slide_count,
    )


def _read_manifest(path: Path) -> dict[str, Any]:
    if not path.is_file():
        return {}
    try:
        value = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        return {}
    return value if isinstance(value, dict) else {}


def _nonempty(path: Path) -> bool:
    return path.is_file() and path.stat().st_size > 0


def _valid_pdf(path: Path) -> bool:
    if not _nonempty(path):
        return False
    try:
        from PyPDF2 import PdfReader

        return len(PdfReader(str(path)).pages) > 0
    except Exception:
        return False


def _valid_pptx(path: Path) -> bool:
    if not _nonempty(path):
        return False
    try:
        from pptx import Presentation

        return len(Presentation(str(path)).slides) > 0
    except Exception:
        return False


def _atomic_write(path: Path, content: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    temporary = path.with_suffix(path.suffix + ".tmp")
    temporary.write_text(content, encoding="utf-8")
    temporary.replace(path)


def _remove_legacy_frontend_layout(chapter_path: Path) -> None:
    """Remove schema-v2 HTML/runtime locations after the nested bundle is valid."""
    (chapter_path / "slides.html").unlink(missing_ok=True)
    legacy_runtime = chapter_path / "frontend-assets"
    if legacy_runtime.is_symlink():
        legacy_runtime.unlink()
    elif legacy_runtime.exists():
        shutil.rmtree(legacy_runtime)
