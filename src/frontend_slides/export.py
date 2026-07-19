from __future__ import annotations

import base64
import tempfile
from pathlib import Path

from .errors import FrontendSlidesError


WIDTH = 1920
HEIGHT = 1080


def export_html_deck(
    html_path: Path,
    *,
    pdf_path: Path | None,
    pptx_path: Path | None,
) -> None:
    """Capture the deck once and create whichever static exports are requested."""
    if pdf_path is None and pptx_path is None:
        return
    if not html_path.is_file():
        raise FrontendSlidesError(f"Presentation HTML not found: {html_path}")
    with tempfile.TemporaryDirectory(prefix="instructional-slides-export-") as tmp:
        screenshot_dir = Path(tmp) / "screenshots"
        screenshot_dir.mkdir()
        screenshots = capture_slide_screenshots(html_path, screenshot_dir)
        if pdf_path is not None:
            _write_pdf(screenshots, pdf_path)
        if pptx_path is not None:
            _write_pptx(screenshots, pptx_path)


def capture_slide_screenshots(html_path: Path, screenshot_dir: Path) -> list[Path]:
    try:
        from playwright.sync_api import Error as PlaywrightError
        from playwright.sync_api import sync_playwright
    except Exception as exc:
        raise FrontendSlidesError(f"Playwright is unavailable for export: {exc}") from exc

    screenshots: list[Path] = []
    try:
        with sync_playwright() as playwright:
            browser = playwright.chromium.launch()
            page = browser.new_page(viewport={"width": WIDTH, "height": HEIGHT})
            page.goto(html_path.resolve().as_uri(), wait_until="load")
            try:
                page.wait_for_function("window.__slidesFitted === true", timeout=30000)
            except Exception:
                pass
            slide_count = page.locator(".slide").count()
            if slide_count < 1:
                raise FrontendSlidesError("No .slide elements were found in slides.html.")
            for index in range(slide_count):
                page.evaluate(_SHOW_SLIDE_JS, index)
                page.wait_for_timeout(120)
                output = screenshot_dir / f"slide-{index + 1:03d}.png"
                page.screenshot(path=str(output), full_page=False)
                screenshots.append(output)
            browser.close()
    except FrontendSlidesError:
        raise
    except PlaywrightError as exc:
        raise FrontendSlidesError(f"Playwright Chromium export failed: {exc}") from exc
    except Exception as exc:
        raise FrontendSlidesError(f"HTML slide capture failed: {exc}") from exc
    return screenshots


def _write_pdf(images: list[Path], output: Path) -> None:
    try:
        from playwright.sync_api import Error as PlaywrightError
        from playwright.sync_api import sync_playwright
    except Exception as exc:
        raise FrontendSlidesError(f"Playwright is unavailable for PDF export: {exc}") from exc

    pages = []
    for image in images:
        encoded = base64.b64encode(image.read_bytes()).decode("ascii")
        pages.append(
            '<section class="page">'
            f'<img src="data:image/png;base64,{encoded}" alt="">'
            "</section>"
        )
    document = f"""<!doctype html>
<html><head><meta charset="utf-8"><style>
* {{ box-sizing:border-box; margin:0; padding:0; }}
@page {{ size:{WIDTH}px {HEIGHT}px; margin:0; }}
.page {{ width:{WIDTH}px; height:{HEIGHT}px; page-break-after:always; overflow:hidden; }}
.page:last-child {{ page-break-after:auto; }}
img {{ display:block; width:100%; height:100%; object-fit:contain; }}
</style></head><body>{''.join(pages)}</body></html>"""
    output.parent.mkdir(parents=True, exist_ok=True)
    try:
        with sync_playwright() as playwright:
            browser = playwright.chromium.launch()
            page = browser.new_page()
            page.set_content(document, wait_until="load")
            page.pdf(
                path=str(output),
                width=f"{WIDTH}px",
                height=f"{HEIGHT}px",
                print_background=True,
                margin={"top": "0", "right": "0", "bottom": "0", "left": "0"},
            )
            browser.close()
    except PlaywrightError as exc:
        raise FrontendSlidesError(f"Playwright PDF export failed: {exc}") from exc
    except Exception as exc:
        raise FrontendSlidesError(f"HTML PDF export failed: {exc}") from exc


def _write_pptx(images: list[Path], output: Path) -> None:
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except Exception as exc:
        raise FrontendSlidesError(f"python-pptx is unavailable: {exc}") from exc

    try:
        presentation = Presentation()
        presentation.slide_width = Inches(13.333333)
        presentation.slide_height = Inches(7.5)
        while presentation.slides:
            relation_id = presentation.slides._sldIdLst[0].rId
            presentation.part.drop_rel(relation_id)
            del presentation.slides._sldIdLst[0]
        layout = presentation.slide_layouts[6]
        for image in images:
            slide = presentation.slides.add_slide(layout)
            slide.shapes.add_picture(
                str(image),
                0,
                0,
                width=presentation.slide_width,
                height=presentation.slide_height,
            )
        output.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(output)
    except Exception as exc:
        raise FrontendSlidesError(f"Static PPTX export failed: {exc}") from exc


_SHOW_SLIDE_JS = """
(index) => {
    const slides = Array.from(document.querySelectorAll('.slide'));
    if (window.presentation && typeof window.presentation.showSlide === 'function') {
        window.presentation.showSlide(index);
    }
    slides.forEach((slide, current) => {
        slide.classList.toggle('active', current === index);
        slide.classList.toggle('visible', current === index);
        slide.style.opacity = current === index ? '1' : '0';
        slide.style.visibility = current === index ? 'visible' : 'hidden';
        slide.style.pointerEvents = current === index ? 'auto' : 'none';
    });
    if (window.__fitAllSlides) window.__fitAllSlides();
    const selected = slides[index];
    if (selected) {
        selected.querySelectorAll('.reveal').forEach((node) => {
            node.style.opacity = '1';
            node.style.transform = 'none';
            node.style.visibility = 'visible';
        });
    }
}
"""
